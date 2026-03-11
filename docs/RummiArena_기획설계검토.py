#!/usr/bin/env python3
"""RummiArena 기획/설계 검토 파워포인트 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ── 색상 팔레트 ──
C_BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)   # 진한 남색 배경
C_BG_CARD   = RGBColor(0x27, 0x27, 0x44)   # 카드 배경
C_PRIMARY   = RGBColor(0x00, 0x96, 0xD6)   # 파란색
C_SECONDARY = RGBColor(0x00, 0xB8, 0x8D)   # 초록/민트
C_ACCENT    = RGBColor(0xFF, 0x6B, 0x35)   # 주황
C_ACCENT2   = RGBColor(0xE0, 0x4F, 0x5F)   # 빨강
C_ACCENT3   = RGBColor(0x9B, 0x59, 0xB6)   # 보라
C_YELLOW    = RGBColor(0xF3, 0xC6, 0x23)   # 노랑
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
C_DIMMED    = RGBColor(0x99, 0x99, 0x99)
C_LINE      = RGBColor(0x44, 0x44, 0x66)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── 헬퍼 함수 ──
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_box(slide, left, top, width, height, fill_color, text="", font_size=10,
            font_color=C_WHITE, bold=False, border_color=None, alignment=PP_ALIGN.CENTER,
            border_width=Pt(1.5)):
    shape = add_rect(slide, left, top, width, height, fill_color, border_color, border_width)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    # 세로 가운데
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_left = Inches(0.08)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.05)
    shape.text_frame.margin_bottom = Inches(0.05)
    return shape

def add_multiline_box(slide, left, top, width, height, fill_color, lines,
                      font_size=10, font_color=C_WHITE, bold_first=True,
                      border_color=None, alignment=PP_ALIGN.CENTER):
    shape = add_rect(slide, left, top, width, height, fill_color, border_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    shape.text_frame.margin_left = Inches(0.1)
    shape.text_frame.margin_right = Inches(0.1)
    shape.text_frame.margin_top = Inches(0.06)
    shape.text_frame.margin_bottom = Inches(0.06)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        if bold_first and i == 0:
            run.font.bold = True
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=C_LIGHT, width=Pt(1.5)):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # head
    connector.end_x = x2
    connector.end_y = y2
    return connector

def add_line(slide, x1, y1, x2, y2, color=C_LINE, width=Pt(1)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_text(slide, left, top, width, height, text, font_size=10,
             font_color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_title_bar(slide, title, subtitle=""):
    # 상단 바
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_BG_DARK)
    add_rect(slide, Inches(0), Inches(0.85), W, Inches(0.04), C_PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.12), Inches(10), Inches(0.5),
             title, font_size=24, font_color=C_WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.5), Inches(10), Inches(0.35),
                 subtitle, font_size=12, font_color=C_DIMMED)
    # 페이지 번호 영역
    add_text(slide, Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.5),
             "RummiArena", font_size=11, font_color=C_PRIMARY, bold=True,
             alignment=PP_ALIGN.RIGHT)

def add_section_label(slide, left, top, text, color=C_PRIMARY):
    add_text(slide, left, top, Inches(4), Inches(0.35),
             text, font_size=13, font_color=color, bold=True)

def inch(v):
    return Inches(v)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)

# 중앙 타이틀 그룹
add_text(slide, inch(0), inch(1.5), W, inch(0.8),
         "RummiArena", font_size=48, font_color=C_PRIMARY, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, inch(0), inch(2.3), W, inch(0.6),
         "멀티 LLM 전략 실험 플랫폼", font_size=28, font_color=C_WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, inch(0), inch(3.2), W, inch(0.5),
         "루미큐브 보드게임 기반  |  Human + AI 혼합 2~4인 실시간 대전", font_size=16,
         font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)

# 하단 라인 + 정보
add_rect(slide, inch(4), inch(4.2), inch(5.333), Pt(2), C_PRIMARY)
add_text(slide, inch(0), inch(4.5), W, inch(0.4),
         "기획/설계 검토 발표자료", font_size=18, font_color=C_YELLOW,
         alignment=PP_ALIGN.CENTER)
add_text(slide, inch(0), inch(5.3), W, inch(0.35),
         "2026-03-11  |  Phase 1 Sprint 0  |  작성: 애벌레",
         font_size=13, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)

# 하단 기술 태그
tags = ["Next.js", "NestJS/Go", "PostgreSQL", "Redis", "OpenAI", "Claude", "DeepSeek",
        "Ollama", "K8s", "ArgoCD", "Helm", "Istio"]
tag_y = inch(6.3)
tag_start = inch(2.3)
for i, tag in enumerate(tags):
    tx = tag_start + inch(i * 0.74)
    add_box(slide, tx, tag_y, inch(0.7), inch(0.28), C_BG_CARD, tag,
            font_size=7, font_color=C_LIGHT, border_color=C_LINE)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: 목차
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "목차", "Table of Contents")

items = [
    ("01", "프로젝트 개요", "목적, 범위, 이해관계자, 성공 기준"),
    ("02", "목표 시스템 구성도", "전체 서비스 구성 및 데이터 흐름"),
    ("03", "논리 아키텍처 구성도", "계층별 컴포넌트 분리 및 인터페이스"),
    ("04", "물리 아키텍처 구성도", "Kubernetes 배포 토폴로지 및 리소스"),
    ("05", "네트워크 구성도", "Ingress 라우팅, 서비스 통신, 외부 연동"),
    ("06", "데이터 아키텍처", "PostgreSQL 스키마, Redis 구조, ER 다이어그램"),
    ("07", "AI Adapter 설계", "멀티 LLM 인터페이스, 캐릭터 시스템, 비용 제어"),
    ("08", "게임 세션 관리", "생명주기, 턴 관리, 장애 복구"),
    ("09", "WBS / 프로젝트 일정", "Sprint-Phase 매핑, 의존관계, 마일스톤"),
    ("10", "핵심 설계 원칙 & 제약", "LLM 신뢰 금지, Stateless, 하드웨어 제약"),
]
for i, (num, title, desc) in enumerate(items):
    row_y = inch(1.3) + inch(i * 0.56)
    # 번호 원
    add_box(slide, inch(1.5), row_y, inch(0.45), inch(0.4), C_PRIMARY, num,
            font_size=14, font_color=C_WHITE, bold=True)
    add_text(slide, inch(2.1), row_y - inch(0.02), inch(4), inch(0.3),
             title, font_size=16, font_color=C_WHITE, bold=True)
    add_text(slide, inch(2.1), row_y + inch(0.25), inch(8), inch(0.25),
             desc, font_size=11, font_color=C_DIMMED)
    # 우측 라인
    add_line(slide, inch(2.1), row_y + inch(0.48), inch(11.5), row_y + inch(0.48), C_LINE, Pt(0.5))


# ══════════════════════════════════════════════════════════════
# SLIDE 3: 프로젝트 개요
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "01. 프로젝트 개요", "Project Overview")

# 좌측: 프로젝트 정보 카드
add_section_label(slide, inch(0.6), inch(1.1), "프로젝트 정보")
info_items = [
    ("프로젝트명", "RummiArena - 멀티 LLM 전략 실험 플랫폼"),
    ("유형", "내부 AI 실험 프로젝트 (외부 서비스 수준 설계)"),
    ("기간", "2026-03-08 ~ 2026-08-15 (약 23주)"),
    ("Sprint", "Sprint 0~9 (2주 주기) + 운영"),
    ("인원", "1인 (PM/설계/개발/운영: 애벌레)"),
    ("저장소", "github.com/k82022603/RummiArena"),
]
for i, (k, v) in enumerate(info_items):
    row_y = inch(1.5) + inch(i * 0.38)
    add_text(slide, inch(0.8), row_y, inch(1.5), inch(0.3),
             k, font_size=10, font_color=C_PRIMARY, bold=True)
    add_text(slide, inch(2.4), row_y, inch(4), inch(0.3),
             v, font_size=10, font_color=C_WHITE)

# 우측 상단: 핵심 목표
add_section_label(slide, inch(7), inch(1.1), "핵심 목표")
goals = [
    ("멀티 LLM 전략 비교", "OpenAI, Claude, DeepSeek, LLaMA\n게임 전략 실험/비교 분석", C_PRIMARY),
    ("풀스택 플랫폼 엔지니어링", "K8s, GitOps, DevSecOps\n전체 사이클 실습", C_SECONDARY),
    ("실시간 멀티플레이", "WebSocket 기반\n2~4인 Human+AI 대전", C_ACCENT),
    ("SaaS 수준 아키텍처", "내부 실험이지만\n외부 공개 가능한 설계", C_ACCENT3),
]
for i, (title, desc, color) in enumerate(goals):
    gx = inch(7) + inch((i % 2) * 3.1)
    gy = inch(1.5) + inch((i // 2) * 1.25)
    add_rect(slide, gx, gy, inch(2.9), inch(1.1), C_BG_CARD, color, Pt(2))
    add_text(slide, gx + inch(0.15), gy + inch(0.1), inch(2.6), inch(0.3),
             title, font_size=12, font_color=color, bold=True)
    add_text(slide, gx + inch(0.15), gy + inch(0.45), inch(2.6), inch(0.6),
             desc, font_size=9, font_color=C_LIGHT)

# 하단: 성공 기준
add_section_label(slide, inch(0.6), inch(4.1), "성공 기준")
criteria = [
    ("게임 완주율", "95%+", C_SECONDARY),
    ("AI 대전", "100판+", C_PRIMARY),
    ("LLM 동시참가", "3개 모델+", C_ACCENT),
    ("CI 파이프라인", "5분 이내", C_ACCENT3),
    ("Pod 복구", "30초 이내", C_ACCENT2),
    ("코드 커버리지", "60%+", C_YELLOW),
    ("보안 취약점", "Critical 0", C_SECONDARY),
    ("AI 턴 응답", "10초 이내", C_PRIMARY),
]
for i, (label, value, color) in enumerate(criteria):
    cx = inch(0.6) + inch(i * 1.55)
    add_rect(slide, cx, inch(4.5), inch(1.4), inch(0.9), C_BG_CARD, color, Pt(1.5))
    add_text(slide, cx + inch(0.1), inch(4.55), inch(1.2), inch(0.3),
             label, font_size=9, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx + inch(0.1), inch(4.85), inch(1.2), inch(0.4),
             value, font_size=18, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)

# 범위
add_section_label(slide, inch(0.6), inch(5.6), "범위 (Scope)")
scope_in = "In: 게임 엔진 / 연습 모드(Stage 1~6) / 멀티플레이 / OAuth / AI 4종 / 캐릭터 시스템 / ELO / 관리자 / K8s / GitOps / DevSecOps / 카카오톡 알림"
scope_out = "Out: 모바일 앱 / 결제 / 100+ 동시 사용자"
add_text(slide, inch(0.8), inch(5.95), inch(11), inch(0.3),
         scope_in, font_size=9, font_color=C_LIGHT)
add_text(slide, inch(0.8), inch(6.3), inch(11), inch(0.3),
         scope_out, font_size=9, font_color=C_ACCENT2)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: 목표 시스템 구성도
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "02. 목표 시스템 구성도", "Target System Architecture")

# ── 사용자 영역 ──
add_section_label(slide, inch(0.5), inch(1.05), "사용자", C_YELLOW)
add_box(slide, inch(0.5), inch(1.45), inch(1.4), inch(0.7), C_BG_CARD,
        "Browser\n(Human Player)", font_size=10, border_color=C_YELLOW)
add_box(slide, inch(0.5), inch(2.35), inch(1.4), inch(0.7), C_BG_CARD,
        "Admin\n(관리자)", font_size=10, border_color=C_YELLOW)

# ── Ingress ──
add_section_label(slide, inch(2.5), inch(1.05), "Ingress (NGINX)", C_LIGHT)
add_box(slide, inch(2.5), inch(1.45), inch(1.6), inch(1.6), C_BG_CARD,
        "NGINX\nIngress\nController\n\nTLS 종단\n라우팅", font_size=9, border_color=C_LIGHT)

# 화살표: 사용자 → Ingress
add_line(slide, inch(1.9), inch(1.8), inch(2.5), inch(1.8), C_YELLOW, Pt(2))
add_line(slide, inch(1.9), inch(2.7), inch(2.5), inch(2.4), C_YELLOW, Pt(2))

# ── Frontend ──
add_box(slide, inch(4.6), inch(1.2), inch(1.6), inch(0.65), C_PRIMARY,
        "Frontend\n(Next.js :3000)", font_size=10, bold=True)
add_box(slide, inch(4.6), inch(2.05), inch(1.6), inch(0.65), C_PRIMARY,
        "Admin Panel\n(Next.js :3001)", font_size=10, bold=True)

# Ingress → Frontend
add_line(slide, inch(4.1), inch(1.65), inch(4.6), inch(1.5), C_PRIMARY, Pt(2))
add_line(slide, inch(4.1), inch(2.4), inch(4.6), inch(2.35), C_PRIMARY, Pt(2))

# ── Game Server ──
gs_x, gs_y = inch(6.8), inch(1.0)
add_rect(slide, gs_x, gs_y, inch(2.4), inch(2.8), C_BG_CARD, C_SECONDARY, Pt(2))
add_text(slide, gs_x + inch(0.1), gs_y + inch(0.05), inch(2.2), inch(0.3),
         "Game Server (:8080)", font_size=12, font_color=C_SECONDARY, bold=True,
         alignment=PP_ALIGN.CENTER)

add_box(slide, gs_x + inch(0.15), gs_y + inch(0.45), inch(2.1), inch(0.5), RGBColor(0x1E, 0x5F, 0x4E),
        "WebSocket + REST API", font_size=10, border_color=C_SECONDARY)
add_box(slide, gs_x + inch(0.15), gs_y + inch(1.1), inch(2.1), inch(0.5), RGBColor(0x1E, 0x5F, 0x4E),
        "Game Engine\n(규칙 검증)", font_size=10, border_color=C_SECONDARY)
add_box(slide, gs_x + inch(0.15), gs_y + inch(1.75), inch(2.1), inch(0.5), RGBColor(0x1E, 0x5F, 0x4E),
        "Session Manager\n(세션/턴 관리)", font_size=10, border_color=C_SECONDARY)

# Frontend → Game Server
add_line(slide, inch(6.2), inch(1.5), inch(6.8), inch(1.5), C_PRIMARY, Pt(2))
add_text(slide, inch(6.2), inch(1.2), inch(0.8), inch(0.25),
         "WS/REST", font_size=8, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)
add_line(slide, inch(6.2), inch(2.35), inch(6.8), inch(2.35), C_PRIMARY, Pt(2))
add_text(slide, inch(6.2), inch(2.55), inch(0.8), inch(0.25),
         "REST", font_size=8, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)

# ── Data Layer ──
add_section_label(slide, inch(6.8), inch(4.1), "데이터 계층", C_ACCENT)
add_box(slide, inch(6.8), inch(4.5), inch(1.5), inch(0.8), RGBColor(0x8B, 0x20, 0x20),
        "Redis 7\n(게임 상태/캐시)\n:6379", font_size=9, border_color=C_ACCENT2)
add_box(slide, inch(8.5), inch(4.5), inch(1.6), inch(0.8), RGBColor(0x1A, 0x3A, 0x6B),
        "PostgreSQL 16\n(영속 데이터)\n:5432", font_size=9, border_color=C_PRIMARY)

# Game Server → Data
add_line(slide, inch(8.0), inch(3.8), inch(7.5), inch(4.5), C_ACCENT2, Pt(2))
add_line(slide, inch(8.2), inch(3.8), inch(9.3), inch(4.5), C_PRIMARY, Pt(2))

# ── AI Adapter ──
ai_x, ai_y = inch(10.0), inch(1.0)
add_rect(slide, ai_x, ai_y, inch(2.8), inch(2.8), C_BG_CARD, C_ACCENT3, Pt(2))
add_text(slide, ai_x + inch(0.1), ai_y + inch(0.05), inch(2.6), inch(0.3),
         "AI Adapter (:8081)", font_size=12, font_color=C_ACCENT3, bold=True,
         alignment=PP_ALIGN.CENTER)

add_box(slide, ai_x + inch(0.1), ai_y + inch(0.45), inch(1.25), inch(0.45),
        RGBColor(0x3D, 0x2A, 0x5C), "PromptBuilder", font_size=9, border_color=C_ACCENT3)
add_box(slide, ai_x + inch(1.45), ai_y + inch(0.45), inch(1.25), inch(0.45),
        RGBColor(0x3D, 0x2A, 0x5C), "ResponseParser", font_size=9, border_color=C_ACCENT3)
add_box(slide, ai_x + inch(0.1), ai_y + inch(1.05), inch(1.25), inch(0.45),
        RGBColor(0x3D, 0x2A, 0x5C), "RetryHandler", font_size=9, border_color=C_ACCENT3)
add_box(slide, ai_x + inch(1.45), ai_y + inch(1.05), inch(1.25), inch(0.45),
        RGBColor(0x3D, 0x2A, 0x5C), "MetricsCollector", font_size=9, border_color=C_ACCENT3)

# Adapters
adapters = [("OpenAI", C_SECONDARY), ("Claude", C_ACCENT), ("DeepSeek", C_PRIMARY), ("Ollama", C_YELLOW)]
for i, (name, c) in enumerate(adapters):
    ax = ai_x + inch(0.1) + inch(i * 0.68)
    add_box(slide, ax, ai_y + inch(1.65), inch(0.62), inch(0.9), RGBColor(0x3D, 0x2A, 0x5C),
            name + "\nAdapter", font_size=7, border_color=c)

# Game Server → AI Adapter
add_line(slide, inch(9.2), inch(2.0), inch(10.0), inch(2.0), C_ACCENT3, Pt(2))
add_text(slide, inch(9.2), inch(1.72), inch(0.9), inch(0.25),
         "gRPC/REST", font_size=8, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)

# ── External LLM APIs ──
add_section_label(slide, inch(10.0), inch(4.1), "외부 LLM APIs", C_ACCENT)
llms = [("OpenAI API", C_SECONDARY), ("Claude API", C_ACCENT), ("DeepSeek API", C_PRIMARY)]
for i, (name, c) in enumerate(llms):
    add_box(slide, inch(10.0) + inch(i * 1.0), inch(4.5), inch(0.9), inch(0.6),
            C_BG_CARD, name, font_size=8, border_color=c)
add_box(slide, inch(10.0), inch(5.3), inch(0.9), inch(0.6), C_BG_CARD,
        "Ollama\n(Local)", font_size=8, border_color=C_YELLOW)
# AI Adapter → LLMs
add_line(slide, inch(11.4), inch(3.8), inch(11.0), inch(4.5), C_ACCENT3, Pt(1.5))

# ── 외부 시스템 ──
add_section_label(slide, inch(0.5), inch(4.1), "외부 시스템", C_ACCENT)
add_box(slide, inch(0.5), inch(4.5), inch(1.4), inch(0.6), C_BG_CARD,
        "Google OAuth 2.0", font_size=9, border_color=C_SECONDARY)
add_box(slide, inch(0.5), inch(5.3), inch(1.4), inch(0.6), C_BG_CARD,
        "Kakao Message API", font_size=9, border_color=C_YELLOW)
# 외부 → Game Server
add_line(slide, inch(1.9), inch(4.8), inch(6.8), inch(3.8), C_SECONDARY, Pt(1))
add_line(slide, inch(1.9), inch(5.6), inch(6.8), inch(3.8), C_YELLOW, Pt(1))

# 범례
add_text(slide, inch(3), inch(6.5), inch(8), inch(0.3),
         "실선: 데이터 흐름  |  색상: 서비스 영역 구분  |  포트는 K8s Service 기준",
         font_size=9, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: 논리 아키텍처 구성도
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "03. 논리 아키텍처 구성도", "Logical Architecture")

layers = [
    ("Presentation Layer", C_PRIMARY, inch(1.1), [
        ("Browser (Next.js)", "게임 UI\ndnd-kit 드래그&드롭\nWebSocket 클라이언트"),
        ("Admin Panel (Next.js)", "관리자 대시보드\n게임 모니터링\nAI 통계"),
    ]),
    ("Application Layer", C_SECONDARY, inch(2.65), [
        ("Auth Module", "Google OAuth\nJWT 발급/검증\nRBAC"),
        ("Game Module", "Room 관리\n턴 오케스트레이션\n세션 생명주기"),
        ("Game Engine", "타일 유효성 검증\n규칙 판정\n승리 조건"),
        ("AI Module", "Adapter 선택\n프롬프트 빌드\n응답 파싱/재시도"),
    ]),
    ("Domain Layer", C_ACCENT, inch(4.2), [
        ("Tile Domain", "타일 인코딩\n{Color}{Num}{Set}\nR7a, JK1"),
        ("Game Domain", "게임 상태 Enum\nWAITING→PLAYING\n→FINISHED"),
        ("Player Domain", "Human/AI 플레이어\nELO 레이팅\n캐릭터 시스템"),
        ("Practice Domain", "연습 모드\nStage 1~6\n목표 판정"),
    ]),
    ("Infrastructure Layer", C_ACCENT2, inch(5.65), [
        ("Redis Client", "게임 상태 R/W\nTTL 관리\nPub/Sub"),
        ("PostgreSQL Client", "유저/전적 CRUD\nELO 이력\n이벤트 로그"),
        ("LLM Client", "OpenAI / Claude\nDeepSeek / Ollama\nHTTP/gRPC"),
        ("WebSocket Server", "연결 관리\n이벤트 브로드캐스트\n재연결 처리"),
    ]),
]

for layer_name, layer_color, ly, components in layers:
    # 레이어 배경
    add_rect(slide, inch(0.4), ly, inch(12.5), inch(1.35), C_BG_CARD, layer_color, Pt(1.5))
    add_text(slide, inch(0.5), ly + inch(0.02), inch(3), inch(0.3),
             layer_name, font_size=11, font_color=layer_color, bold=True)

    comp_count = len(components)
    comp_w = inch(2.8) if comp_count <= 2 else inch(2.85)
    gap = inch(0.2)
    total_w = comp_count * comp_w + (comp_count - 1) * gap
    start_x = inch(0.4) + (inch(12.5) - total_w) / 2

    for j, (comp_name, comp_desc) in enumerate(components):
        cx = start_x + j * (comp_w + gap)
        add_multiline_box(slide, cx, ly + inch(0.35), comp_w, inch(0.9),
                          RGBColor(0x1E, 0x1E, 0x36), [comp_name, comp_desc],
                          font_size=9, bold_first=True, border_color=layer_color)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: 물리 아키텍처 구성도
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "04. 물리 아키텍처 구성도", "Physical Architecture (Kubernetes)")

# ── Host Machine ──
add_rect(slide, inch(0.3), inch(1.05), inch(12.7), inch(6.1), C_BG_CARD, C_DIMMED, Pt(1))
add_text(slide, inch(0.5), inch(1.1), inch(5), inch(0.3),
         "LG Gram 15Z90R  |  Windows 11 + WSL2  |  RAM 16GB  |  i7-1360P",
         font_size=10, font_color=C_DIMMED)

# ── WSL2 ──
add_rect(slide, inch(0.5), inch(1.5), inch(12.3), inch(5.5), RGBColor(0x20, 0x20, 0x3A), C_LINE, Pt(1))
add_text(slide, inch(0.7), inch(1.55), inch(3), inch(0.25),
         "WSL2 (Ubuntu)  |  Memory: 10~14GB", font_size=9, font_color=C_LINE)

# ── Docker Desktop K8s ──
add_rect(slide, inch(0.7), inch(1.9), inch(8.5), inch(4.95), RGBColor(0x1A, 0x1A, 0x30), C_PRIMARY, Pt(1.5))
add_text(slide, inch(0.9), inch(1.95), inch(4), inch(0.25),
         "Docker Desktop Kubernetes (docker-desktop)", font_size=10, font_color=C_PRIMARY, bold=True)

# Namespace: rummikub
add_rect(slide, inch(0.9), inch(2.3), inch(8.1), inch(4.4), RGBColor(0x18, 0x18, 0x2E), C_SECONDARY, Pt(1))
add_text(slide, inch(1.1), inch(2.35), inch(3), inch(0.25),
         "Namespace: rummikub", font_size=9, font_color=C_SECONDARY, bold=True)

# Deployments
deploys = [
    ("frontend", "Next.js\n:3000\nreplicas: 1", "128Mi/256Mi", C_PRIMARY),
    ("game-server", "NestJS/Go\n:8080\nreplicas: 1", "256Mi/512Mi", C_SECONDARY),
    ("ai-adapter", "NestJS/Go\n:8081\nreplicas: 1", "256Mi/512Mi", C_ACCENT3),
    ("admin", "Next.js\n:3001\nreplicas: 1", "128Mi/256Mi", C_PRIMARY),
]
for i, (name, desc, mem, color) in enumerate(deploys):
    dx = inch(1.1) + inch(i * 2.0)
    add_rect(slide, dx, inch(2.75), inch(1.8), inch(1.7), RGBColor(0x22, 0x22, 0x3E), color, Pt(1.5))
    add_text(slide, dx + inch(0.05), inch(2.8), inch(1.7), inch(0.25),
             "Deployment", font_size=7, font_color=C_DIMMED)
    add_text(slide, dx + inch(0.05), inch(3.0), inch(1.7), inch(0.25),
             name, font_size=11, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, dx + inch(0.05), inch(3.3), inch(1.7), inch(0.7),
             desc, font_size=8, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text(slide, dx + inch(0.05), inch(4.15), inch(1.7), inch(0.2),
             mem, font_size=7, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)

# StatefulSets
sts = [
    ("Redis 7", "StatefulSet\n:6379\n100Mi/200Mi", C_ACCENT2),
    ("PostgreSQL 16", "StatefulSet + PVC\n:5432\n256Mi/512Mi", RGBColor(0x33, 0x66, 0xCC)),
]
for i, (name, desc, color) in enumerate(sts):
    dx = inch(1.1) + inch(i * 2.0)
    add_rect(slide, dx, inch(4.65), inch(1.8), inch(1.2), RGBColor(0x22, 0x22, 0x3E), color, Pt(1.5))
    add_text(slide, dx + inch(0.05), inch(4.7), inch(1.7), inch(0.25),
             "StatefulSet", font_size=7, font_color=C_DIMMED)
    add_text(slide, dx + inch(0.05), inch(4.9), inch(1.7), inch(0.25),
             name, font_size=11, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, dx + inch(0.05), inch(5.2), inch(1.7), inch(0.5),
             desc, font_size=8, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)

# Optional Pod (Ollama)
add_rect(slide, inch(5.1), inch(4.65), inch(1.8), inch(1.2), RGBColor(0x22, 0x22, 0x3E), C_YELLOW, Pt(1))
add_text(slide, inch(5.15), inch(4.7), inch(1.7), inch(0.25),
         "K8s 밖 실행 (선택)", font_size=7, font_color=C_DIMMED)
add_text(slide, inch(5.15), inch(4.9), inch(1.7), inch(0.25),
         "Ollama", font_size=11, font_color=C_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, inch(5.15), inch(5.2), inch(1.7), inch(0.5),
         "LLaMA 1B~3B\n:11434\n~3GB RAM", font_size=8, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)

# ── 우측: 교대 실행 전략 ──
add_rect(slide, inch(9.5), inch(1.9), inch(3.3), inch(4.95), C_BG_CARD, C_ACCENT, Pt(1.5))
add_text(slide, inch(9.7), inch(1.95), inch(3), inch(0.3),
         "교대 실행 전략 (리소스 제약)", font_size=11, font_color=C_ACCENT, bold=True)

modes = [
    ("개발 모드", "앱 + Redis + PG", "~6GB", C_SECONDARY),
    ("CI/CD 모드", "ArgoCD + Runner", "~6GB", C_PRIMARY),
    ("품질 모드", "SonarQube (단독)", "~4GB", C_ACCENT3),
    ("AI 실험 모드", "Ollama + 1~3B", "~5GB", C_YELLOW),
]
for i, (mode, svc, ram, color) in enumerate(modes):
    my = inch(2.4) + inch(i * 1.05)
    add_rect(slide, inch(9.7), my, inch(2.9), inch(0.9), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1))
    add_text(slide, inch(9.85), my + inch(0.05), inch(2.6), inch(0.25),
             mode, font_size=10, font_color=color, bold=True)
    add_text(slide, inch(9.85), my + inch(0.3), inch(2.6), inch(0.25),
             svc, font_size=9, font_color=C_LIGHT)
    add_text(slide, inch(9.85), my + inch(0.55), inch(2.6), inch(0.25),
             "RAM " + ram, font_size=10, font_color=C_WHITE, bold=True)

# 총 예상
add_text(slide, inch(9.7), inch(6.55), inch(3), inch(0.3),
         "* 모든 서비스 동시 실행 불가", font_size=8, font_color=C_ACCENT2)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: 네트워크 구성도
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "05. 네트워크 구성도", "Network Architecture")

# ── 외부 영역 ──
add_section_label(slide, inch(0.5), inch(1.1), "External (Internet)", C_YELLOW)
ext_items = [
    ("Browser", inch(0.5), C_YELLOW),
    ("Google OAuth", inch(2.5), C_SECONDARY),
    ("Kakao API", inch(4.5), C_YELLOW),
    ("LLM APIs\n(OpenAI, Claude\nDeepSeek)", inch(6.5), C_ACCENT),
]
for name, ex, color in ext_items:
    add_box(slide, ex, inch(1.5), inch(1.6), inch(0.7), C_BG_CARD, name,
            font_size=9, border_color=color)

# ── Ingress 영역 ──
add_rect(slide, inch(0.3), inch(2.6), inch(12.5), inch(1.2), C_BG_CARD, C_LIGHT, Pt(1))
add_text(slide, inch(0.5), inch(2.65), inch(4), inch(0.25),
         "NGINX Ingress Controller  (TLS Termination, self-signed cert)", font_size=10,
         font_color=C_LIGHT, bold=True)

routes = [
    ("/", "frontend:3000", C_PRIMARY),
    ("/api/*", "game-server:8080", C_SECONDARY),
    ("/ws", "game-server:8080\n(WebSocket Upgrade)", C_SECONDARY),
    ("/admin/*", "admin:3001", C_PRIMARY),
]
for i, (path, target, color) in enumerate(routes):
    rx = inch(0.5) + inch(i * 3.1)
    add_box(slide, rx, inch(3.0), inch(1.0), inch(0.55), RGBColor(0x1E, 0x1E, 0x36),
            path, font_size=12, font_color=color, bold=True)
    add_text(slide, rx + inch(1.1), inch(3.1), inch(1.8), inch(0.45),
             "-> " + target, font_size=9, font_color=C_LIGHT)

# Browser → Ingress
add_line(slide, inch(1.3), inch(2.2), inch(1.3), inch(2.6), C_YELLOW, Pt(2))
add_text(slide, inch(0.5), inch(2.3), inch(1.5), inch(0.2),
         "HTTPS / WSS", font_size=8, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)

# ── ClusterIP Services ──
add_rect(slide, inch(0.3), inch(4.1), inch(12.5), inch(1.5), RGBColor(0x18, 0x18, 0x2E), C_SECONDARY, Pt(1))
add_text(slide, inch(0.5), inch(4.15), inch(5), inch(0.25),
         "ClusterIP Services  (K8s Internal Network, Namespace: rummikub)", font_size=10,
         font_color=C_SECONDARY, bold=True)

svcs = [
    ("frontend\nClusterIP\n:3000", C_PRIMARY),
    ("game-server\nClusterIP\n:8080", C_SECONDARY),
    ("ai-adapter\nClusterIP\n:8081", C_ACCENT3),
    ("admin\nClusterIP\n:3001", C_PRIMARY),
    ("redis\nClusterIP\n:6379", C_ACCENT2),
    ("postgres\nClusterIP\n:5432", RGBColor(0x33, 0x66, 0xCC)),
]
for i, (name, color) in enumerate(svcs):
    sx = inch(0.5) + inch(i * 2.1)
    add_box(slide, sx, inch(4.5), inch(1.8), inch(0.9), RGBColor(0x1E, 0x1E, 0x36),
            name, font_size=9, border_color=color)

# ── 서비스 간 통신 매트릭스 ──
add_section_label(slide, inch(0.5), inch(5.85), "서비스 간 통신", C_ACCENT)
comm_items = [
    "game-server -> redis:6379 (TCP, 게임 상태 R/W)",
    "game-server -> postgres:5432 (TCP, 유저/전적/로그)",
    "game-server -> ai-adapter:8081 (HTTP/gRPC, AI 행동 요청)",
    "ai-adapter -> LLM APIs (HTTPS, 외부 모델 호출)",
    "ai-adapter -> ollama:11434 (HTTP, 로컬 모델 - K8s 밖)",
]
for i, item in enumerate(comm_items):
    add_text(slide, inch(0.7) + inch((i // 3) * 6), inch(6.15) + inch((i % 3) * 0.3),
             inch(6), inch(0.3), item, font_size=9, font_color=C_LIGHT)

# Phase 5 참고
add_text(slide, inch(8), inch(5.85), inch(4.5), inch(0.25),
         "Phase 5 (Sprint 9): Istio mTLS, VirtualService, DestinationRule 추가",
         font_size=9, font_color=C_ACCENT3)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: 데이터 아키텍처
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "06. 데이터 아키텍처", "Data Architecture")

# ── 좌측: PostgreSQL ──
add_section_label(slide, inch(0.5), inch(1.1), "PostgreSQL 16 (영속 데이터)", C_PRIMARY)

pg_tables = [
    ("users", "id, email, display_name\nelo_rating, role, is_blocked", C_PRIMARY),
    ("games", "id, room_code, status\ngame_mode, settings(JSONB)", C_SECONDARY),
    ("game_players", "player_type, seat_order\nai_persona, ai_difficulty\nai_psychology_level", C_ACCENT),
    ("ai_call_logs", "player_type, model_name\nai_persona, latency_ms\nprompt/completion_tokens", C_ACCENT3),
    ("game_events", "turn_number, event_type\nevent_data(JSONB)", C_ACCENT2),
    ("elo_history", "rating_before/after\nrating_delta, k_factor", C_YELLOW),
    ("practice_sessions", "stage(1~6), status\nobjectives, result(JSONB)", C_SECONDARY),
]
for i, (name, fields, color) in enumerate(pg_tables):
    col = i % 2
    row = i // 2
    tx = inch(0.5) + inch(col * 3.2)
    ty = inch(1.5) + inch(row * 1.35)
    add_rect(slide, tx, ty, inch(3.0), inch(1.2), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1.5))
    add_text(slide, tx + inch(0.1), ty + inch(0.05), inch(2.8), inch(0.25),
             name, font_size=11, font_color=color, bold=True)
    add_text(slide, tx + inch(0.1), ty + inch(0.35), inch(2.8), inch(0.8),
             fields, font_size=8, font_color=C_LIGHT)

# system_config
add_rect(slide, inch(0.5), inch(6.2), inch(6.2), inch(0.45), RGBColor(0x1E, 0x1E, 0x36), C_DIMMED, Pt(1))
add_text(slide, inch(0.6), inch(6.22), inch(6), inch(0.35),
         "system_config: turn_timeout_sec=60, ai_max_retries=3, ai_timeout_ms=10000, max_rooms=10",
         font_size=8, font_color=C_DIMMED)

# 관계 표시
relations = [
    ("users 1--* game_players", inch(0.5)),
    ("games 1--* game_players, game_events", inch(1.0)),
    ("game_players 1--* ai_call_logs", inch(1.5)),
    ("users 1--* elo_history, practice_sessions", inch(2.0)),
]

# ── 우측: Redis ──
add_section_label(slide, inch(7.2), inch(1.1), "Redis 7 (실시간 상태)", C_ACCENT2)

redis_keys = [
    ("game:{id}:state", "Hash: status, currentTurn,\ncurrentPlayer, drawPileCount\nTTL: 7200s (매 턴 갱신)", C_ACCENT2),
    ("game:{id}:player:{seat}:tiles", 'List: ["R1a","B5a","JK1",...]\nTTL: 7200s', C_ACCENT),
    ("game:{id}:drawpile", 'List: 셔플된 타일 목록\nTTL: 7200s', C_PRIMARY),
    ("game:{id}:stalemate", "Hash: consecutiveDrawCount,\nlastPlaceTurn\nTTL: 7200s", C_ACCENT3),
    ("game:{id}:timer", "String: 만료 Unix timestamp\nTTL: 30~120s (턴 설정)", C_YELLOW),
    ("session:{id}", "Hash: userId, gameId\nTTL: 1800s", C_SECONDARY),
    ("quota:daily:{date}", "Hash: totalCalls, totalCost\n모델별 호출수/비용\nTTL: 172800s", C_ACCENT2),
]
for i, (key, desc, color) in enumerate(redis_keys):
    ty = inch(1.5) + inch(i * 0.75)
    add_rect(slide, inch(7.2), ty, inch(5.5), inch(0.65), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1))
    add_text(slide, inch(7.35), ty + inch(0.03), inch(2.5), inch(0.25),
             key, font_size=9, font_color=color, bold=True)
    add_text(slide, inch(9.8), ty + inch(0.03), inch(2.8), inch(0.6),
             desc, font_size=8, font_color=C_LIGHT)

# 타일 인코딩 범례
add_rect(slide, inch(7.2), inch(6.85), inch(5.5), inch(0.35), C_BG_CARD, C_LINE, Pt(1))
add_text(slide, inch(7.35), inch(6.87), inch(5.2), inch(0.3),
         "타일 인코딩: {Color}{Number}{Set}  |  R7a=빨강7(a), B13b=파랑13(b), JK1=조커1  |  총 106장",
         font_size=8, font_color=C_LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: AI Adapter 설계
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "07. AI Adapter 설계", "AI Adapter Architecture & Character System")

# ── 좌측: 인터페이스 흐름 ──
add_section_label(slide, inch(0.5), inch(1.1), "요청 흐름 (LLM 신뢰 금지)", C_ACCENT3)

flow_items = [
    ("Game Server", "턴 요청", C_SECONDARY, inch(0.5)),
    ("PromptBuilder", "게임상태→프롬프트", C_ACCENT3, inch(2.5)),
    ("LLM API", "모델 호출", C_ACCENT, inch(4.5)),
    ("ResponseParser", "JSON 파싱", C_ACCENT3, inch(6.5)),
    ("Game Engine", "유효성 검증", C_SECONDARY, inch(8.5)),
]
for name, desc, color, fx in flow_items:
    add_box(slide, fx, inch(1.5), inch(1.6), inch(0.7), RGBColor(0x1E, 0x1E, 0x36),
            name + "\n" + desc, font_size=9, border_color=color)
# 화살표
for i in range(len(flow_items) - 1):
    x1 = flow_items[i][3] + inch(1.6)
    x2 = flow_items[i+1][3]
    add_line(slide, x1, inch(1.85), x2, inch(1.85), C_LIGHT, Pt(1.5))

# 재시도 루프
add_text(slide, inch(6.5), inch(2.3), inch(4.5), inch(0.25),
         "무효 -> 재요청(max 3회) -> 전부 실패 -> 강제 드로우",
         font_size=9, font_color=C_ACCENT2)
add_box(slide, inch(10.5), inch(1.5), inch(1.5), inch(0.7), RGBColor(0x5A, 0x1A, 0x1A),
        "실패 시\n강제 드로우", font_size=9, border_color=C_ACCENT2)
add_line(slide, inch(10.1), inch(1.85), inch(10.5), inch(1.85), C_ACCENT2, Pt(1.5))

# ── 캐릭터 시스템 ──
add_section_label(slide, inch(0.5), inch(2.7), "AI 캐릭터 시스템 (6 x 3 x 4)", C_YELLOW)

chars = [
    ("Rookie\n(루키)", "순진, 실수", "하수", C_SECONDARY),
    ("Calculator\n(칼큘레이터)", "논리적, 차분", "중수", C_PRIMARY),
    ("Shark\n(샤크)", "공격적, 압박형", "고수", C_ACCENT2),
    ("Fox\n(폭스)", "교활, 전략형", "고수", C_ACCENT),
    ("Wall\n(월)", "수비적, 끈질김", "중수", RGBColor(0x33, 0x66, 0xCC)),
    ("Wildcard\n(와일드카드)", "예측불가", "중수", C_ACCENT3),
]
for i, (name, trait, diff, color) in enumerate(chars):
    cx = inch(0.5) + inch(i * 2.1)
    add_rect(slide, cx, inch(3.1), inch(1.9), inch(1.3), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1.5))
    add_text(slide, cx + inch(0.05), inch(3.15), inch(1.8), inch(0.5),
             name, font_size=10, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, cx + inch(0.05), inch(3.65), inch(1.8), inch(0.25),
             trait, font_size=8, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)
    add_box(slide, cx + inch(0.5), inch(4.0), inch(0.9), inch(0.25), color,
            diff, font_size=8, font_color=C_WHITE, bold=True)

# ── 난이도별 차이 ──
add_section_label(slide, inch(0.5), inch(4.6), "난이도별 정보 제한", C_ACCENT)

diff_data = [
    ("하수 (Beginner)", "경량 모델\ngpt-4o-mini, llama3.2:1b\n내 타일+테이블만\n재배치 불가, 심리전 0", C_SECONDARY),
    ("중수 (Intermediate)", "중급 모델\ngpt-4o-mini, deepseek\n+ 상대 남은 수\n기본 재배치, 심리전 1~2", C_PRIMARY),
    ("고수 (Expert)", "최상위 모델\ngpt-4o, claude-sonnet\n+ 히스토리+미출현 타일\n적극 재배치, 심리전 2~3", C_ACCENT2),
]
for i, (level, desc, color) in enumerate(diff_data):
    dx = inch(0.5) + inch(i * 4.2)
    add_rect(slide, dx, inch(5.0), inch(3.9), inch(1.55), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1.5))
    add_text(slide, dx + inch(0.1), inch(5.05), inch(3.7), inch(0.3),
             level, font_size=11, font_color=color, bold=True)
    add_text(slide, dx + inch(0.1), inch(5.35), inch(3.7), inch(1.1),
             desc, font_size=9, font_color=C_LIGHT)

# 비용 제어
add_text(slide, inch(0.5), inch(6.7), inch(12), inch(0.3),
         "비용 제어: 사용자당 500회/일, 게임당 200회, 일일 $10 한도  |  프롬프트 토큰 예산: 2,000 이내  |  히스토리: 최근 N턴 제한 (난이도별 차등)",
         font_size=9, font_color=C_DIMMED)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: 게임 세션 관리
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "08. 게임 세션 관리", "Game Session Lifecycle & Turn Management")

# ── 생명주기 ──
add_section_label(slide, inch(0.5), inch(1.1), "세션 생명주기", C_SECONDARY)

states = [
    ("WAITING", "Room 생성\n플레이어 대기", C_YELLOW, inch(1.0)),
    ("PLAYING", "게임 진행\n턴 순환", C_SECONDARY, inch(4.0)),
    ("FINISHED", "정상 종료\n승자 확정", C_PRIMARY, inch(7.0)),
    ("CANCELLED", "비정상 종료\n인원 부족/강제", C_ACCENT2, inch(7.0)),
]
# WAITING
add_box(slide, inch(1.0), inch(1.5), inch(2.0), inch(0.8), RGBColor(0x1E, 0x1E, 0x36),
        "WAITING\nRoom 생성, 플레이어 대기", font_size=10, border_color=C_YELLOW)
# PLAYING
add_box(slide, inch(4.0), inch(1.5), inch(2.0), inch(0.8), RGBColor(0x1E, 0x1E, 0x36),
        "PLAYING\n게임 진행, 턴 순환", font_size=10, border_color=C_SECONDARY)
# FINISHED
add_box(slide, inch(7.0), inch(1.3), inch(2.0), inch(0.65), RGBColor(0x1E, 0x1E, 0x36),
        "FINISHED\n정상 종료", font_size=10, border_color=C_PRIMARY)
# CANCELLED
add_box(slide, inch(7.0), inch(2.1), inch(2.0), inch(0.65), RGBColor(0x1E, 0x1E, 0x36),
        "CANCELLED\n비정상 종료", font_size=10, border_color=C_ACCENT2)

# 전이 화살표
add_line(slide, inch(3.0), inch(1.9), inch(4.0), inch(1.9), C_LIGHT, Pt(2))
add_text(slide, inch(3.0), inch(1.6), inch(1.0), inch(0.2), "호스트 시작\n(2명+)",
         font_size=7, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)
add_line(slide, inch(6.0), inch(1.65), inch(7.0), inch(1.65), C_LIGHT, Pt(2))
add_text(slide, inch(6.0), inch(1.4), inch(1.0), inch(0.2), "승자 확정",
         font_size=7, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)
add_line(slide, inch(6.0), inch(2.1), inch(7.0), inch(2.35), C_ACCENT2, Pt(2))

# 종료 처리 요약
add_rect(slide, inch(9.5), inch(1.3), inch(3.3), inch(1.5), C_BG_CARD, C_PRIMARY, Pt(1))
add_text(slide, inch(9.65), inch(1.35), inch(3), inch(0.25),
         "종료 처리", font_size=10, font_color=C_PRIMARY, bold=True)
end_steps = "1. 승자 판정\n2. 점수 계산 (조커=30점)\n3. ELO 레이팅 업데이트\n4. PostgreSQL 저장\n5. game:ended 이벤트\n6. Redis TTL 600s로 단축"
add_text(slide, inch(9.65), inch(1.65), inch(3), inch(1.1),
         end_steps, font_size=8, font_color=C_LIGHT)

# ── 턴 관리 ──
add_section_label(slide, inch(0.5), inch(3.1), "턴 관리 흐름", C_ACCENT)

# Human 턴
add_rect(slide, inch(0.5), inch(3.5), inch(5.8), inch(1.8), C_BG_CARD, C_YELLOW, Pt(1.5))
add_text(slide, inch(0.65), inch(3.55), inch(3), inch(0.25),
         "Human 턴", font_size=11, font_color=C_YELLOW, bold=True)

h_steps = [
    ("turn:start", "이벤트 전송", C_YELLOW),
    ("타이머 시작", "30~120초", C_ACCENT),
    ("드래그&드롭", "타일 배치", C_PRIMARY),
    ("turn:confirm", "검증 요청", C_SECONDARY),
]
for i, (s, d, c) in enumerate(h_steps):
    hx = inch(0.7) + inch(i * 1.4)
    add_box(slide, hx, inch(3.95), inch(1.2), inch(0.55), RGBColor(0x1E, 0x1E, 0x36),
            s + "\n" + d, font_size=8, border_color=c)
    if i < len(h_steps) - 1:
        add_line(slide, hx + inch(1.2), inch(4.2), hx + inch(1.4), inch(4.2), C_LIGHT, Pt(1))

add_box(slide, inch(0.7), inch(4.7), inch(2.5), inch(0.45), RGBColor(0x3A, 0x1A, 0x1A),
        "타임아웃 -> 자동 드로우", font_size=9, border_color=C_ACCENT2)
add_box(slide, inch(3.5), inch(4.7), inch(2.5), inch(0.45), RGBColor(0x1A, 0x3A, 0x1A),
        "유효 -> 적용 -> 다음 턴", font_size=9, border_color=C_SECONDARY)

# AI 턴
add_rect(slide, inch(6.8), inch(3.5), inch(5.8), inch(1.8), C_BG_CARD, C_ACCENT3, Pt(1.5))
add_text(slide, inch(6.95), inch(3.55), inch(3), inch(0.25),
         "AI 턴", font_size=11, font_color=C_ACCENT3, bold=True)

a_steps = [
    ("ai:thinking", "사고 중 표시", C_ACCENT3),
    ("AI Adapter", "LLM 호출", C_ACCENT),
    ("Engine 검증", "유효성 확인", C_SECONDARY),
    ("turn:action", "결과 브로드캐스트", C_PRIMARY),
]
for i, (s, d, c) in enumerate(a_steps):
    ax = inch(7.0) + inch(i * 1.4)
    add_box(slide, ax, inch(3.95), inch(1.2), inch(0.55), RGBColor(0x1E, 0x1E, 0x36),
            s + "\n" + d, font_size=8, border_color=c)
    if i < len(a_steps) - 1:
        add_line(slide, ax + inch(1.2), inch(4.2), ax + inch(1.4), inch(4.2), C_LIGHT, Pt(1))

add_box(slide, inch(7.0), inch(4.7), inch(2.5), inch(0.45), RGBColor(0x3A, 0x1A, 0x1A),
        "3회 실패 -> 강제 드로우", font_size=9, border_color=C_ACCENT2)
add_box(slide, inch(9.8), inch(4.7), inch(2.5), inch(0.45), RGBColor(0x1A, 0x3A, 0x1A),
        "유효 -> 적용 -> 다음 턴", font_size=9, border_color=C_SECONDARY)

# ── 장애 복구 ──
add_section_label(slide, inch(0.5), inch(5.6), "연결 끊김 / 장애 복구", C_ACCENT2)

recovery = [
    ("Human 끊김", "30초 재연결 유예\n-> 자동 드로우\n-> 3턴 부재시 제외\n-> 2명 미만시 종료", C_YELLOW),
    ("AI 장애", "재시도 max 3회\n-> 강제 드로우\n-> 5회 연속시\n   AI 비활성화", C_ACCENT3),
    ("서버 재시작", "Redis 스캔\n-> 타이머 복구\n-> 30초 재연결 대기\n-> 복구 or 종료", C_SECONDARY),
    ("ELO 규칙", "정상 종료만 적용\nCANCELLED 미적용\n연습 모드 미적용\nK=32/40/24", C_PRIMARY),
]
for i, (title, desc, color) in enumerate(recovery):
    rx = inch(0.5) + inch(i * 3.2)
    add_rect(slide, rx, inch(5.95), inch(3.0), inch(1.2), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1.5))
    add_text(slide, rx + inch(0.1), inch(6.0), inch(2.8), inch(0.25),
             title, font_size=10, font_color=color, bold=True)
    add_text(slide, rx + inch(0.1), inch(6.3), inch(2.8), inch(0.8),
             desc, font_size=8, font_color=C_LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: WBS / 프로젝트 일정
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "09. WBS / 프로젝트 일정", "Work Breakdown Structure & Timeline")

# ── 타임라인 바 ──
add_section_label(slide, inch(0.5), inch(1.1), "프로젝트 타임라인 (23주)", C_PRIMARY)

sprints = [
    ("S0", "기획\n환경", "03/08\n~03/28", C_DIMMED, 3),
    ("S1", "게임\n엔진", "03/29\n~04/11", C_SECONDARY, 2),
    ("S2", "백엔드\nAPI", "04/12\n~04/25", C_SECONDARY, 2),
    ("S3", "프론트\n연습모드", "04/26\n~05/09", C_SECONDARY, 2),
    ("S4", "AI\nAdapter", "05/10\n~05/23", C_ACCENT3, 2),
    ("S5", "멀티\n플레이", "05/24\n~06/06", C_ACCENT3, 2),
    ("S6", "관리자\n랭킹", "06/07\n~06/20", C_ACCENT, 2),
    ("S7", "Observ\nability", "06/21\n~07/04", C_ACCENT2, 2),
    ("S8", "보안\n고도화", "07/05\n~07/18", C_ACCENT2, 2),
    ("S9", "Service\nMesh", "07/19\n~08/01", C_ACCENT2, 2),
    ("OP", "운영\n실험", "08/02\n~08/15", C_YELLOW, 2),
]

total_weeks = sum(s[4] for s in sprints)
bar_start = inch(0.5)
bar_width = inch(12.0)
bar_y = inch(1.5)
bar_h = inch(1.5)

x_offset = bar_start
for name, desc, dates, color, weeks in sprints:
    w = bar_width * weeks / total_weeks
    add_rect(slide, x_offset, bar_y, w - Pt(2), bar_h, RGBColor(0x1E, 0x1E, 0x36), color, Pt(1.5))
    add_text(slide, x_offset + Pt(4), bar_y + inch(0.02), w - Pt(8), inch(0.22),
             name, font_size=9, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x_offset + Pt(4), bar_y + inch(0.28), w - Pt(8), inch(0.5),
             desc, font_size=7, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text(slide, x_offset + Pt(4), bar_y + inch(0.9), w - Pt(8), inch(0.5),
             dates, font_size=6, font_color=C_DIMMED, alignment=PP_ALIGN.CENTER)
    x_offset += w

# Phase 라벨
phases = [
    ("Phase 1", inch(0.5), 3, C_DIMMED),
    ("Phase 2 (MVP)", inch(0.5) + bar_width * 3 / total_weeks, 6, C_SECONDARY),
    ("Phase 3", inch(0.5) + bar_width * 9 / total_weeks, 4, C_ACCENT3),
    ("Phase 4", inch(0.5) + bar_width * 13 / total_weeks, 2, C_ACCENT),
    ("Phase 5", inch(0.5) + bar_width * 15 / total_weeks, 6, C_ACCENT2),
    ("Phase 6", inch(0.5) + bar_width * 21 / total_weeks, 2, C_YELLOW),
]
for pname, px, pw, pc in phases:
    pwidth = bar_width * pw / total_weeks
    add_text(slide, px, inch(3.05), pwidth, inch(0.25),
             pname, font_size=9, font_color=pc, bold=True, alignment=PP_ALIGN.CENTER)

# ── Sprint 상세 내용 ──
add_section_label(slide, inch(0.5), inch(3.5), "Sprint별 주요 산출물", C_SECONDARY)

sprint_details = [
    ("Sprint 0", "기획 문서 5종\n설계 문서 5종\nK8s/ArgoCD/GitLab\nBackend 기술 결정"),
    ("Sprint 1~3", "게임 엔진 (검증/턴)\nREST+WebSocket API\nRedis/PG 연동\nNext.js UI, 연습 모드"),
    ("Sprint 4~5", "AI Adapter 4종\n캐릭터 시스템\n프롬프트 설계\nHuman+AI 통합 대전"),
    ("Sprint 6", "관리자 대시보드\nELO 랭킹 시스템\n카카오톡 알림\nAI 통계"),
    ("Sprint 7~9", "Prometheus/Grafana\nSonarQube/Trivy/ZAP\nIstio/Kiali/Jaeger\nk6 부하 테스트"),
    ("운영", "AI vs AI 토너먼트\n모델별 전략 분석\n프롬프트 최적화\n운영 가이드"),
]
for i, (sprint, detail) in enumerate(sprint_details):
    sx = inch(0.5) + inch(i * 2.1)
    add_rect(slide, sx, inch(3.9), inch(2.0), inch(1.8), RGBColor(0x1E, 0x1E, 0x36), C_LINE, Pt(1))
    add_text(slide, sx + inch(0.1), inch(3.95), inch(1.8), inch(0.25),
             sprint, font_size=10, font_color=C_WHITE, bold=True)
    add_text(slide, sx + inch(0.1), inch(4.25), inch(1.8), inch(1.4),
             detail, font_size=8, font_color=C_LIGHT)

# ── 의존관계 ──
add_section_label(slide, inch(0.5), inch(5.9), "Sprint 의존 관계", C_ACCENT)
dep_text = ("S0 -> S1 (Backend 결정)  |  S1 -> S2 (엔진 로직)  |  S1,S2 -> S3 (API+엔진)  |  "
            "S2 -> S4 (API 연동)  |  S3,S4 -> S5 (통합)  |  S5 -> S6 -> S7 -> S8 -> S9 -> 운영")
add_text(slide, inch(0.7), inch(6.2), inch(12), inch(0.4),
         dep_text, font_size=9, font_color=C_LIGHT)

# 핵심 마일스톤
add_section_label(slide, inch(0.5), inch(6.6), "핵심 마일스톤", C_YELLOW)
ms = [
    ("2026-03-28", "기획 완료, 인프라 구축"),
    ("2026-05-09", "MVP (게임 엔진+API+UI)"),
    ("2026-06-06", "AI 대전 완성"),
    ("2026-08-01", "DevSecOps 고도화"),
    ("2026-08-15", "프로젝트 완료"),
]
for i, (date, desc) in enumerate(ms):
    mx = inch(0.5) + inch(i * 2.6)
    add_text(slide, mx, inch(6.9), inch(1.0), inch(0.25),
             date, font_size=9, font_color=C_YELLOW, bold=True)
    add_text(slide, mx + inch(1.0), inch(6.9), inch(1.5), inch(0.25),
             desc, font_size=9, font_color=C_LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: 핵심 설계 원칙 & 제약
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_title_bar(slide, "10. 핵심 설계 원칙 & 제약 사항", "Key Design Principles & Constraints")

# ── 설계 원칙 ──
add_section_label(slide, inch(0.5), inch(1.1), "핵심 설계 원칙", C_PRIMARY)

principles = [
    ("LLM 신뢰 금지", "LLM은 행동 '제안'만 수행\nGame Engine이 모든 유효성 검증\n무효 시 재요청(3회) → 강제 드로우",
     C_ACCENT2, "01"),
    ("AI Adapter 분리", "Game Engine은 특정 LLM에 의존 안 함\n공통 인터페이스(MoveRequest/Response)\n모델 교체 가능, Istio 트래픽 분배",
     C_ACCENT3, "02"),
    ("Stateless Server", "게임 상태는 Redis에 저장\nPod 재시작 시 게임 유지\nRedis Pub/Sub로 수평 확장 대비",
     C_SECONDARY, "03"),
    ("GitOps + DevSecOps", "소스 repo ↔ GitOps repo 분리\nArgoCD Helm 기반 자동 배포\nSonarQube + Trivy + ZAP 보안 게이트",
     C_PRIMARY, "04"),
]
for i, (title, desc, color, num) in enumerate(principles):
    px = inch(0.5) + inch(i * 3.15)
    add_rect(slide, px, inch(1.5), inch(3.0), inch(2.2), C_BG_CARD, color, Pt(2))
    add_box(slide, px + inch(0.1), inch(1.55), inch(0.4), inch(0.35), color,
            num, font_size=14, font_color=C_WHITE, bold=True)
    add_text(slide, px + inch(0.6), inch(1.55), inch(2.3), inch(0.35),
             title, font_size=13, font_color=color, bold=True)
    add_text(slide, px + inch(0.15), inch(2.0), inch(2.7), inch(1.5),
             desc, font_size=9, font_color=C_LIGHT)

# ── 기술 스택 ──
add_section_label(slide, inch(0.5), inch(3.9), "기술 스택 요약", C_SECONDARY)

stack_groups = [
    ("Frontend", "Next.js\nTailwindCSS\nFramer Motion\ndnd-kit", C_PRIMARY),
    ("Backend", "NestJS or Go\n(Sprint 0 내 결정)\nWebSocket\nREST API", C_SECONDARY),
    ("Database", "PostgreSQL 16\nRedis 7\n8 테이블\n7 Redis 키 패턴", C_ACCENT),
    ("AI", "OpenAI API\nClaude API\nDeepSeek API\nOllama (LLaMA)", C_ACCENT3),
    ("Infra", "Docker Desktop K8s\nHelm 3\nArgoCD\nNGINX Ingress", C_ACCENT2),
    ("Quality", "SonarQube\nTrivy\nOWASP ZAP\nk6", C_YELLOW),
    ("Auth", "Google OAuth 2.0\nJWT + RBAC\nADMIN / USER", C_PRIMARY),
    ("Notification", "Kakao\nMessage API\n게임 결과/알림", C_ACCENT),
]
for i, (name, items, color) in enumerate(stack_groups):
    sx = inch(0.5) + inch(i * 1.6)
    add_rect(slide, sx, inch(4.3), inch(1.5), inch(1.5), RGBColor(0x1E, 0x1E, 0x36), color, Pt(1))
    add_text(slide, sx + inch(0.05), inch(4.32), inch(1.4), inch(0.25),
             name, font_size=10, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, sx + inch(0.05), inch(4.6), inch(1.4), inch(1.1),
             items, font_size=8, font_color=C_LIGHT, alignment=PP_ALIGN.CENTER)

# ── 제약 사항 ──
add_section_label(slide, inch(0.5), inch(6.0), "주요 제약 사항", C_ACCENT2)
constraints = [
    "RAM 16GB: 교대 실행 전략 필수, 모든 서비스 동시 실행 불가",
    "1인 개발: 23주 10 Sprint, 2주 주기",
    "GPU 없음: Ollama 3B 이하 모델만 로컬 실행",
    "LLM 비용: 일일 $10 한도, 호출 수 제한",
]
for i, c in enumerate(constraints):
    add_text(slide, inch(0.7) + inch((i // 2) * 6.3), inch(6.35) + inch((i % 2) * 0.3),
             inch(6), inch(0.3), "- " + c, font_size=9, font_color=C_LIGHT)


# ══════════════════════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════════════════════
output_path = "/mnt/d/Users/KTDS/Documents/06.과제/RummiArena/docs/RummiArena_기획설계검토.pptx"
prs.save(output_path)
print(f"파워포인트 저장 완료: {output_path}")
print(f"총 {len(prs.slides)}장 슬라이드 생성")
