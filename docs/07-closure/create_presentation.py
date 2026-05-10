"""
RummiArena 프로젝트 종료 보고서 PowerPoint 생성 스크립트
Tech Innovation 테마 (primary=#2E7D32, secondary=#81C784, accent=#FF9800)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ─── 테마 ───────────────────────────────────────────────────────────────────
THEME = {
    'primary':   RGBColor(46, 125, 50),      # #2E7D32
    'secondary': RGBColor(129, 199, 132),    # #81C784
    'accent':    RGBColor(255, 152, 0),      # #FF9800
    'dark':      RGBColor(27, 94, 32),       # #1B5E20
    'light':     RGBColor(232, 245, 233),    # #E8F5E9
    'text':      RGBColor(33, 33, 33),       # #212121
    'white':     RGBColor(255, 255, 255),
    'gray':      RGBColor(120, 120, 120),
    'red':       RGBColor(198, 40, 40),
    'blue':      RGBColor(21, 101, 192),
    'yellow':    RGBColor(245, 127, 23),
    'desc_bg':   RGBColor(226, 240, 217),    # 우측 설명 박스 배경
}

FONT = "맑은 고딕"

# ─── 공통 헬퍼 ──────────────────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    s.line.width = line_w
    return s

def add_round(slide, l, t, w, h, fill, line=None, line_w=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    s.line.width = line_w
    return s

def add_tb(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tb

def add_tb_multiline(slide, l, t, w, h, lines, size, color, bold=False, space_after=Pt(6)):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.space_after = space_after
    return tb

def add_arrow(slide, x1, y1, x2, y2, color=None, width=Pt(2.5)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color or THEME['accent']
    c.line.width = width
    c.line.end_arrow_type = 2
    return c

def slide_header(slide, title, subtitle=None):
    """공통 헤더: 진한 녹색 바 + 제목"""
    add_rect(slide, 0, 0, 10, 0.85, THEME['dark'])
    add_tb(slide, 0.35, 0.1, 8.5, 0.65, title, 26, THEME['white'], bold=True)
    if subtitle:
        add_tb(slide, 0.35, 0.55, 9, 0.35, subtitle, 12, THEME['secondary'])
    # 하단 구분선
    add_rect(slide, 0, 7.2, 10, 0.05, THEME['secondary'])
    add_tb(slide, 0.3, 7.25, 9, 0.22, "RummiArena — LLM 전략 실험 플랫폼 | 프로젝트 종료 보고서 2026-05-10",
           9, THEME['gray'])

def desc_box(slide, lines, l=6.65, t=0.95, w=3.15, h=6.15):
    """우측 설명 박스 (연두색 배경)"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = THEME['desc_bg']
    s.line.color.rgb = THEME['desc_bg']
    tb = slide.shapes.add_textbox(Inches(l+0.12), Inches(t+0.12), Inches(w-0.24), Inches(h-0.24))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for i, line in enumerate(lines):
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        if line.startswith("■"):
            p.text = line; p.font.size = Pt(10); p.font.bold = True
            p.font.color.rgb = THEME['dark']; p.font.name = FONT
        else:
            p.text = line; p.font.size = Pt(9.5)
            p.font.color.rgb = THEME['text']; p.font.name = FONT
        p.space_after = Pt(3)

# ─── 슬라이드 1: 표지 ────────────────────────────────────────────────────────
def slide_title(prs):
    slide = new_slide(prs)
    # 배경
    add_rect(slide, 0, 0, 10, 7.5, THEME['dark'])
    # 장식 원
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(4.5), Inches(5), Inches(4))
    s.fill.solid(); s.fill.fore_color.rgb = THEME['primary']
    s.line.color.rgb = THEME['primary']
    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-0.5), Inches(4), Inches(3))
    s2.fill.solid(); s2.fill.fore_color.rgb = THEME['primary']
    s2.line.color.rgb = THEME['primary']
    # 중앙 배너
    add_rect(slide, 0.5, 2.2, 9, 0.08, THEME['secondary'])
    add_rect(slide, 0.5, 4.8, 9, 0.08, THEME['secondary'])
    # 제목
    add_tb(slide, 0.8, 1.2, 8.5, 1.1, "RummiArena", 54, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 0.8, 2.4, 8.5, 0.7, "프로젝트 종료 보고서", 32, THEME['secondary'], align=PP_ALIGN.CENTER)
    # 부제목
    add_tb(slide, 0.8, 3.15, 8.5, 0.55,
           "루미큐브 기반 멀티 LLM 전략 실험 플랫폼",
           20, THEME['light'], align=PP_ALIGN.CENTER)
    # 메타
    add_tb(slide, 0.8, 4.95, 8.5, 0.5,
           "프로젝트 기간: 2026-03-08 ~ 2026-05-10  (63일)",
           16, THEME['secondary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 0.8, 5.5, 8.5, 0.5,
           "팀 구성: 에이전트 13명 | 커밋 750+ | 테스트 2,462건",
           14, THEME['light'], align=PP_ALIGN.CENTER)
    add_tb(slide, 0.8, 6.3, 8.5, 0.5,
           "작성: 애벌레 (프로젝트 오너)  |  2026-05-10",
           13, THEME['gray'], align=PP_ALIGN.CENTER)

# ─── 슬라이드 2: 목차 ────────────────────────────────────────────────────────
def slide_toc(prs):
    slide = new_slide(prs)
    slide_header(slide, "목 차")
    sections = [
        ("1", "프로젝트 개요 및 팀 구성"),
        ("2", "시스템 아키텍처 및 기술 스택"),
        ("3", "개발 여정 — Sprint 1~7 + 핫픽스"),
        ("4", "AI 대전 실험 결과 (LLM 비교)"),
        ("5", "UI 아키텍처 — pendingStore SSOT"),
        ("6", "게임룰 SSOT — 77개 룰 체계"),
        ("7", "품질 · 보안 · CI/CD"),
        ("8", "운영 환경 및 동시접속 분석"),
        ("9", "기술 부채 및 향후 과제"),
        ("10", "프로젝트 최종 성과 및 마무리"),
    ]
    col1 = sections[:5]; col2 = sections[5:]
    for i, (num, text) in enumerate(col1):
        y = 1.1 + i * 1.08
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.35), Inches(y), Inches(0.55), Inches(0.5))
        s.fill.solid(); s.fill.fore_color.rgb = THEME['primary']; s.line.color.rgb = THEME['primary']
        add_tb(slide, 0.35, y-0.02, 0.55, 0.54, num, 16, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, 1.05, y+0.04, 3.9, 0.42, text, 15, THEME['text'])
    for i, (num, text) in enumerate(col2):
        y = 1.1 + i * 1.08
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.15), Inches(y), Inches(0.55), Inches(0.5))
        s.fill.solid(); s.fill.fore_color.rgb = THEME['dark']; s.line.color.rgb = THEME['dark']
        add_tb(slide, 5.15, y-0.02, 0.55, 0.54, num, 16, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, 5.85, y+0.04, 3.9, 0.42, text, 15, THEME['text'])

# ─── 슬라이드 3: 프로젝트 개요 ───────────────────────────────────────────────
def slide_overview(prs):
    slide = new_slide(prs)
    slide_header(slide, "프로젝트 개요", "RummiArena — 루미큐브 기반 멀티 LLM 전략 실험 플랫폼")
    # 3개 카드
    cards = [
        (THEME['primary'], "목적",
         ["루미큐브 게임 환경에서", "GPT, Claude, DeepSeek,", "Ollama(LLaMA) 등 다양한", "LLM의 전략 능력을", "비교·실험하는 플랫폼"]),
        (THEME['dark'], "범위",
         ["Human + AI 혼합 2~4인", "실시간 멀티플레이 대전", "웹 기반 게임 UI + 관리자", "K8s 기반 운영 환경", "63일 / Sprint 7 + 핫픽스"]),
        (THEME['accent'], "성과",
         ["커밋 750+  |  문서 110+", "테스트 2,462건 통과", "AI 최고 place rate 33.3%", "LLaMA 로컬 모델 25.6%", "CI/CD 17/17 ALL GREEN"]),
    ]
    for i, (color, title, items) in enumerate(cards):
        x = 0.3 + i * 3.2
        add_round(slide, x, 0.95, 3.0, 5.9, color, line=color)
        add_tb(slide, x+0.1, 1.0, 2.8, 0.5, title, 17, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, x+0.05, 1.48, 2.9, 0.04, THEME['white'])
        tb = slide.shapes.add_textbox(Inches(x+0.15), Inches(1.6), Inches(2.7), Inches(4.8))
        tf = tb.text_frame; tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "• " + item; p.font.size = Pt(13); p.font.name = FONT
            p.font.color.rgb = THEME['white']; p.space_after = Pt(8)

# ─── 슬라이드 4: 팀 구성 ─────────────────────────────────────────────────────
def slide_team(prs):
    slide = new_slide(prs)
    slide_header(slide, "팀 구성 — 에이전트 13명")
    agents = [
        ("애벌레",          "프로젝트 오너",        THEME['dark']),
        ("pm",              "프로젝트 매니저",       THEME['primary']),
        ("architect",       "시스템 설계",           THEME['primary']),
        ("go-dev",          "게임서버 (Go)",         THEME['blue']),
        ("node-dev",        "AI 어댑터 (NestJS)",    THEME['blue']),
        ("frontend-dev",    "프론트엔드 (Next.js)",  THEME['blue']),
        ("frontend-dev-opus","프론트 페어코딩",      THEME['blue']),
        ("devops",          "K8s/CI/CD",             THEME['accent']),
        ("qa",              "테스트 품질관리",        THEME['accent']),
        ("security",        "보안 엔지니어",          THEME['red']),
        ("ai-engineer",     "LLM 프롬프트",           THEME['primary']),
        ("game-analyst",    "게임룰 SSOT",            THEME['primary']),
        ("designer",        "UI/UX 설계",             THEME['secondary']),
    ]
    cols = 4
    for i, (name, role, color) in enumerate(agents):
        col = i % cols
        row = i // cols
        x = 0.3 + col * 2.42
        y = 1.05 + row * 1.6
        s = add_round(slide, x, y, 2.2, 1.35, color, line=color)
        add_tb(slide, x+0.08, y+0.1, 2.05, 0.5, name, 14, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.08, y+0.62, 2.05, 0.55, role, 11, THEME['white'], align=PP_ALIGN.CENTER)

# ─── 슬라이드 5: 시스템 아키텍처 ────────────────────────────────────────────
def slide_architecture(prs):
    slide = new_slide(prs)
    slide_header(slide, "시스템 아키텍처", "5개 서비스 + Redis + PostgreSQL — Kubernetes 운영")

    main_w = 6.2

    # 사용자 영역
    add_round(slide, 0.3, 0.95, 1.55, 0.75, THEME['secondary'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.3, 1.0, 1.55, 0.65, "👤 사용자\n(Browser)", 11, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)

    add_round(slide, 0.3, 1.85, 1.55, 0.75, THEME['secondary'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.3, 1.9, 1.55, 0.65, "👤 관리자\n(Browser)", 11, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)

    # 화살표 → Frontend
    add_arrow(slide, 1.85, 1.33, 2.3, 1.55, THEME['primary'])
    add_arrow(slide, 1.85, 2.23, 2.3, 2.43, THEME['primary'])

    # Frontend / Admin
    add_round(slide, 2.3, 1.05, 1.6, 0.75, THEME['primary'])
    add_tb(slide, 2.3, 1.1, 1.6, 0.65, "Frontend\n(Next.js)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 2.3, 1.95, 1.6, 0.75, THEME['dark'])
    add_tb(slide, 2.3, 2.0, 1.6, 0.65, "Admin Panel\n(Next.js)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 화살표 → Game Server
    add_arrow(slide, 3.9, 1.43, 4.3, 1.43, THEME['accent'])
    add_arrow(slide, 3.9, 2.33, 4.3, 2.33, THEME['accent'])

    # Game Server
    add_round(slide, 4.3, 0.95, 1.8, 1.5, THEME['dark'])
    add_tb(slide, 4.3, 1.0, 1.8, 0.4, "Game Server", 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 4.3, 1.4, 1.8, 0.95,
           "Go / gin\ngorilla/ws\nGORM\nGame Engine", 10, THEME['light'], align=PP_ALIGN.CENTER)

    # 화살표 → AI Adapter
    add_arrow(slide, 6.1, 1.7, 6.5, 1.7, THEME['accent'])

    # AI Adapter
    add_round(slide, 6.5, 1.0, 1.65, 1.35, THEME['primary'])
    add_tb(slide, 6.5, 1.05, 1.65, 0.4, "AI Adapter", 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 6.5, 1.45, 1.65, 0.8,
           "NestJS\nTypeScript\nRate Limit", 10, THEME['white'], align=PP_ALIGN.CENTER)

    # AI Adapter → LLMs
    llms = [("GPT-5-mini\n(OpenAI)", 0.9), ("DeepSeek\nV4-Pro", 1.65), ("Claude\nSonnet 4", 2.4), ("Ollama\nLLaMA 3B", 3.15)]
    for label, t_off in llms:
        add_round(slide, 8.28, 1.0 + (t_off-0.9)/0.75*0.6 - 0.05, 1.42, 0.6, THEME['accent'], line=THEME['accent'])
        add_tb(slide, 8.28, 1.0 + (t_off-0.9)/0.75*0.6 - 0.02, 1.42, 0.6, label, 9, THEME['white'], align=PP_ALIGN.CENTER)
        add_arrow(slide, 8.15, 1.33 + (t_off-0.9)/0.75*0.6 - 0.02, 8.28, 1.33 + (t_off-0.9)/0.75*0.6 - 0.02, THEME['secondary'], Pt(1.5))

    # 재배치
    for i, (label, y_i) in enumerate([("GPT-5-mini\n(OpenAI)", 0), ("DeepSeek\nV4-Pro", 1), ("Claude\nSonnet 4", 2), ("Ollama\nLLaMA 3B", 3)]):
        yy = 1.0 + i * 0.78
        add_round(slide, 8.25, yy, 1.45, 0.65, THEME['accent'], line=THEME['accent'])
        add_tb(slide, 8.25, yy+0.02, 1.45, 0.62, label, 9, THEME['white'], align=PP_ALIGN.CENTER)
        add_arrow(slide, 8.15, yy+0.32, 8.25, yy+0.32, THEME['secondary'], Pt(1.5))

    # 데이터 저장소 (하단)
    add_round(slide, 0.5, 4.0, 2.4, 0.8, THEME['blue'])
    add_tb(slide, 0.5, 4.05, 2.4, 0.7, "Redis 7\n(게임 상태)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 3.2, 4.0, 2.4, 0.8, THEME['blue'])
    add_tb(slide, 3.2, 4.05, 2.4, 0.7, "PostgreSQL 16\n(영속 데이터)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 5.9, 4.0, 2.0, 0.8, THEME['gray'])
    add_tb(slide, 5.9, 4.05, 2.0, 0.7, "Istio\nService Mesh", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 8.1, 4.0, 1.6, 0.8, THEME['gray'])
    add_tb(slide, 8.1, 4.05, 1.6, 0.7, "ArgoCD\n(GitOps)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # Game Server → DB
    add_arrow(slide, 5.2, 2.45, 5.2, 4.0, THEME['blue'])
    add_arrow(slide, 4.9, 2.45, 1.7, 4.0, THEME['blue'])

    # 레이어 라벨
    add_tb(slide, 0.3, 4.9, 9.4, 0.35,
           "K8s (Docker Desktop)  |  Namespace: rummikub  |  7개 서비스 Running  |  Helm 차트 5개",
           12, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.25, 4.85, 9.5, 0.04, THEME['secondary'])

    # 설명 박스
    desc_box(slide, [
        "■ 핵심 설계 원칙",
        "",
        "1. LLM 신뢰 금지",
        "   LLM 응답은 반드시 Game",
        "   Engine이 유효성 검증.",
        "   Invalid → 재시도 3회",
        "   → 강제 드로우",
        "",
        "2. AI Adapter 분리",
        "   모든 LLM은 동일한",
        "   MoveRequest/Response",
        "   인터페이스로 통신.",
        "   모델 교체 가능.",
        "",
        "3. Stateless 서버",
        "   게임 상태 = Redis.",
        "   Pod 재시작에도 게임",
        "   상태 유지.",
        "",
        "4. GitOps",
        "   GitHub(SSOT) →",
        "   GitLab(mirror) →",
        "   ArgoCD → K8s 배포.",
    ], l=6.55, t=0.92, w=3.2, h=6.15)

# ─── 슬라이드 6: 기술 스택 ───────────────────────────────────────────────────
def slide_tech_stack(prs):
    slide = new_slide(prs)
    slide_header(slide, "기술 스택", "Frontend · Backend · AI · Infra · Quality")
    categories = [
        ("Frontend", THEME['primary'], [
            "Next.js (React 18)",
            "TailwindCSS",
            "dnd-kit (드래그앤드롭)",
            "Framer Motion (애니메이션)",
            "next-auth (Google OAuth)",
        ]),
        ("Backend", THEME['dark'], [
            "Go / gin + gorilla/ws",
            "GORM (PostgreSQL)",
            "NestJS / TypeScript",
            "Redis 7 (Pub/Sub)",
            "JWT 인증 미들웨어",
        ]),
        ("AI / LLM", THEME['accent'], [
            "OpenAI gpt-5-mini",
            "Claude Sonnet 4 (thinking)",
            "DeepSeek V4-Pro (thinking)",
            "Ollama qwen2.5:3b (로컬)",
            "프롬프트 v1~v9 진화",
        ]),
        ("Infra / CI", THEME['blue'], [
            "Docker Desktop K8s",
            "Helm 3 + ArgoCD",
            "GitLab CI (17 stages)",
            "Istio Service Mesh",
            "Trivy + SonarQube",
        ]),
    ]
    for i, (cat, color, items) in enumerate(categories):
        x = 0.3 + i * 2.42
        add_round(slide, x, 0.95, 2.2, 5.95, color, line=color)
        add_tb(slide, x+0.05, 1.0, 2.1, 0.5, cat, 15, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, x+0.1, 1.5, 2.0, 0.04, THEME['white'])
        tb = slide.shapes.add_textbox(Inches(x+0.15), Inches(1.6), Inches(1.9), Inches(5.0))
        tf = tb.text_frame; tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "• " + item; p.font.size = Pt(12); p.font.name = FONT
            p.font.color.rgb = THEME['white']; p.space_after = Pt(10)

# ─── 슬라이드 7: 스프린트 타임라인 ──────────────────────────────────────────
def slide_sprint_timeline(prs):
    slide = new_slide(prs)
    slide_header(slide, "개발 여정 — 63일 스프린트 타임라인")
    sprints = [
        ("S1", "03-08", "03-24", "기반 구축"),
        ("S2", "03-25", "04-04", "AI 캐릭터"),
        ("S3", "04-05", "04-14", "OAuth/WS"),
        ("S4", "04-15", "04-22", "보안 P0"),
        ("S5", "04-23", "04-30", "Rate Limit"),
        ("S6", "05-01", "05-07", "재배치 UI"),
        ("S7", "04-25", "05-04", "마라톤"),
        ("HF", "05-01", "05-10", "핫픽스"),
    ]
    colors = [THEME['secondary'], THEME['secondary'], THEME['primary'], THEME['primary'],
              THEME['dark'], THEME['dark'], THEME['accent'], THEME['red']]

    # 타임라인 축
    add_rect(slide, 0.7, 4.1, 9.0, 0.06, THEME['gray'])
    for i, (sp, s, e, desc) in enumerate(sprints):
        x = 0.7 + i * 1.1
        y = 1.05 if i % 2 == 0 else 2.55
        # 박스
        add_round(slide, x, y, 1.0, 0.95, colors[i], line=colors[i])
        add_tb(slide, x, y+0.04, 1.0, 0.38, sp, 16, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x, y+0.42, 1.0, 0.5, desc, 9, THEME['white'], align=PP_ALIGN.CENTER)
        # 수직선
        if i % 2 == 0:
            add_arrow(slide, x+0.5, y+0.95, x+0.5, 4.1, colors[i], Pt(1.5))
        else:
            add_arrow(slide, x+0.5, 4.1, x+0.5, y+0.95, colors[i], Pt(1.5))
        # 날짜
        add_tb(slide, x-0.0, 4.2, 1.05, 0.3, s, 8, THEME['gray'], align=PP_ALIGN.CENTER)
        add_tb(slide, x-0.0, 4.5, 1.05, 0.3, "~"+e, 8, THEME['gray'], align=PP_ALIGN.CENTER)

    # 성과 요약 (하단)
    milestones = [
        ("750+ 커밋", THEME['primary']),
        ("110+ 문서", THEME['dark']),
        ("2,462 테스트", THEME['accent']),
        ("CI/CD 17/17", THEME['blue']),
        ("LLM 4종 실험", THEME['primary']),
    ]
    for i, (text, color) in enumerate(milestones):
        x = 0.3 + i * 1.85
        add_round(slide, x, 5.0, 1.7, 0.65, color, line=color)
        add_tb(slide, x, 5.05, 1.7, 0.55, text, 13, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── 슬라이드 8: Sprint 1~4 ──────────────────────────────────────────────────
def slide_sprint_1_4(prs):
    slide = new_slide(prs)
    slide_header(slide, "Sprint 1~4 — 기반 구축부터 보안 강화까지")
    items = [
        ("Sprint 1  (03-08~03-24)  기반 구축", THEME['primary'], [
            "Game Engine 핵심 (그룹/런 검증, 초기 용융, 조커)",
            "REST API + WebSocket 구현",
            "K8s 5개 서비스 첫 배포 (Helm)",
            "JWT fail-fast, CVE 7건 수정",
        ]),
        ("Sprint 2  (03-25~04-04)  AI 캐릭터 + ELO", THEME['primary'], [
            "AI 캐릭터 6종 × 3난이도 설계",
            "Turn Orchestrator + ELO 랭킹 시스템",
            "관리자 대시보드 + 연습 모드",
            "첫 LLM 프롬프트 v1 실험",
        ]),
        ("Sprint 3  (04-05~04-14)  OAuth + WS 재연결", THEME['dark'], [
            "Google OAuth 2.0 K8s 통합",
            "WebSocket 재연결 + Redis Session",
            "Redis Timer 기반 턴 타임아웃",
            "Istio Service Mesh 도입 (game-server + ai-adapter)",
        ]),
        ("Sprint 4  (04-15~04-22)  보안 P0 + DeepSeek", THEME['dark'], [
            "보안 5건: SQL injection, XSS, Admin API 인증",
            "WebSocket 인증 검증, dev-login 차단",
            "DeepSeek V4-Pro 통합 (23.1% place rate)",
            "비용 한도: 일일 $20 / 시간당 $5",
        ]),
    ]
    for i, (title, color, bullets) in enumerate(items):
        row = i // 2; col = i % 2
        x = 0.3 + col * 4.85; y = 1.0 + row * 3.15
        add_round(slide, x, y, 4.6, 2.95, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, x, y, 4.6, 0.52, color)
        add_tb(slide, x+0.1, y+0.06, 4.4, 0.42, title, 12, THEME['white'], bold=True)
        tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.62), Inches(4.2), Inches(2.2))
        tf = tb.text_frame; tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "• " + b; p.font.size = Pt(12); p.font.name = FONT
            p.font.color.rgb = THEME['text']; p.space_after = Pt(6)

# ─── 슬라이드 9: Sprint 5~7 + 핫픽스 ────────────────────────────────────────
def slide_sprint_5_7(prs):
    slide = new_slide(prs)
    slide_header(slide, "Sprint 5~7 + 핫픽스 — 고도화·마라톤·실측 완료")
    items = [
        ("Sprint 5  (04-23~04-30)  Rate Limit + CI", THEME['primary'], [
            "Rate Limiting: REST Sliding Window / WS Fixed Window",
            "CI/CD 17/17 ALL GREEN 달성",
            "DeepSeek N=3 실험: 23.1% 확정",
            "플레이테스트 88.6% 통과",
        ]),
        ("Sprint 6  (05-01~05-07)  재배치 UI", THEME['primary'], [
            "dnd-kit 드래그앤드롭 도입",
            "재배치 4유형 UI 구현 (손→보드, 보드→보드 등)",
            "Agent Teams 13명 체제 본격화",
            "핫픽스 4건 처리",
        ]),
        ("Sprint 7  (04-25~05-04)  마라톤 마감", THEME['accent'], [
            "pendingStore SSOT 아키텍처 (38개 파일, 13개 필드 제거)",
            "게임룰 SSOT 71→77개 룰 확정",
            "SEC-A/B/C 완료 (govulncheck 25→0건)",
            "Phase D: 8.5 Story Point × 1.5시간 달성",
        ]),
        ("핫픽스 세션  (05-01~05-10)  실측 완료", THEME['red'], [
            "BUG-CONFIRM-001: confirmBusy 영구 잠금 수정",
            "LLaMA v8(15.8%) → v9(25.6%) place rate 달성",
            "V-22 단독 승리 룰, 로비 ELO API 연동",
            "Jest 659 / AI Adapter 637 PASS 최종 확인",
        ]),
    ]
    for i, (title, color, bullets) in enumerate(items):
        row = i // 2; col = i % 2
        x = 0.3 + col * 4.85; y = 1.0 + row * 3.15
        add_round(slide, x, y, 4.6, 2.95, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, x, y, 4.6, 0.52, color)
        add_tb(slide, x+0.1, y+0.06, 4.4, 0.42, title, 12, THEME['white'], bold=True)
        tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.62), Inches(4.2), Inches(2.2))
        tf = tb.text_frame; tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "• " + b; p.font.size = Pt(12); p.font.name = FONT
            p.font.color.rgb = THEME['text']; p.space_after = Pt(6)

# ─── 슬라이드 10: AI 대전 결과 ───────────────────────────────────────────────
def slide_llm_results(prs):
    slide = new_slide(prs)
    slide_header(slide, "AI 대전 실험 결과 — LLM 4종 비교", "place rate = 턴당 실제 타일 배치 성공 비율")

    # 막대그래프 (도형으로 구현)
    models = [
        ("GPT\ngpt-5-mini",     33.3, THEME['primary'],   "v2 프롬프트\n$0.10/1M"),
        ("DeepSeek\nV4-Pro",    31.9, THEME['dark'],      "thinking 모드\n$0.039/게임"),
        ("Ollama\nLLaMA 3B",    25.6, THEME['accent'],    "로컬 CPU\n무료"),
        ("Claude\nSonnet 4",    20.0, THEME['blue'],      "thinking 모드\n$1.11/게임"),
    ]
    max_h = 3.5
    for i, (name, rate, color, note) in enumerate(models):
        x = 0.55 + i * 2.3
        bar_h = (rate / 35.0) * max_h
        y_top = 4.6 - bar_h
        # 막대
        add_round(slide, x, y_top, 1.8, bar_h, color, line=color)
        # 퍼센트 라벨
        add_tb(slide, x, y_top - 0.38, 1.8, 0.35, f"{rate}%", 20, color, bold=True, align=PP_ALIGN.CENTER)
        # 모델명
        add_tb(slide, x, 4.65, 1.8, 0.55, name, 12, THEME['text'], align=PP_ALIGN.CENTER)
        # 부가정보
        add_tb(slide, x, 5.25, 1.8, 0.5, note, 10, THEME['gray'], align=PP_ALIGN.CENTER)

    # 기준선
    add_rect(slide, 0.4, 4.62, 9.2, 0.05, THEME['gray'])

    # 요약 박스
    add_round(slide, 0.35, 5.85, 9.3, 1.1, THEME['light'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.55, 5.9, 9.0, 0.4,
           "핵심 발견: Thinking 모드가 필수 — V4-Flash(비thinking) 0% vs V4-Pro(thinking) 31.9%",
           13, THEME['dark'], bold=True)
    add_tb(slide, 0.55, 6.3, 9.0, 0.55,
           "비용 대비 성능: DeepSeek $0.039/게임 vs Claude $1.11/게임 (28배 차이, 성능 비슷)  |  LLaMA 로컬 무료 25.6%",
           12, THEME['text'])

    desc_box(slide, [
        "■ 실험 조건",
        "",
        "N=3 self-play",
        "(게임 3회 평균)",
        "",
        "■ 평가 지표",
        "• place rate:",
        "  턴당 배치 성공률",
        "• fallback:",
        "  엔진 거부 후",
        "  강제 드로우",
        "",
        "■ GPT v2 고정 이유",
        "v4 프롬프트 실험 결과",
        "reasoning_tokens",
        "-25% 감소(Cohen d",
        "=-1.46). 내부 RLHF가",
        "외부 지시를 무시.",
        "",
        "■ LLaMA 목표",
        "v10에서 30%+ 도전",
        "(ERR_NO_RACK_TILE",
        "수정 시 가능 전망)",
    ], l=6.55, t=0.92, w=3.2, h=6.15)

# ─── 슬라이드 11: Ollama 진화 ────────────────────────────────────────────────
def slide_ollama(prs):
    slide = new_slide(prs)
    slide_header(slide, "LLaMA(Ollama) 프롬프트 진화 — v6→v9", "3B 파라미터 CPU 추론 모델에게 루미큐브를 가르친 8주")
    versions = [
        ("v6", "단순 질의",     "0%",   "모델이 단 한 번도\n타일을 배치 안 함",    THEME['red']),
        ("v7", "배치 판단 추가", "~5%",  "낮은 성공률\n개선 방향 탐색 중",         THEME['red']),
        ("v8", "사전 계산 전략", "15.8%","가능한 조합 목록화\n→ 선택 방식 도입",   THEME['accent']),
        ("v9", "ERR 방어+조커", "25.6%","ERR_NO_RACK_TILE 방어\n+ 조커 지원 개선", THEME['primary']),
    ]
    for i, (ver, strategy, rate, desc, color) in enumerate(versions):
        x = 0.35 + i * 2.42
        # 버전 원
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.55), Inches(1.05), Inches(1.1), Inches(1.1))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.color.rgb = color
        add_tb(slide, x+0.55, 1.12, 1.1, 0.95, ver, 20, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        # 화살표
        if i < 3:
            add_arrow(slide, x+2.15, 1.6, x+2.42, 1.6, color, Pt(2))
        add_tb(slide, x+0.1, 2.3, 2.1, 0.45, strategy, 12, THEME['text'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.1, 2.75, 2.1, 0.55, rate, 26, color, bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.1, 3.35, 2.1, 0.85, desc, 11, THEME['gray'], align=PP_ALIGN.CENTER)

    # 핵심 변화 설명
    add_round(slide, 0.35, 4.35, 6.1, 2.55, THEME['light'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.55, 4.4, 5.8, 0.4, "v8 → v9 핵심 변화 (+9.8%p 향상)", 14, THEME['dark'], bold=True)
    changes = [
        "• ERR_NO_RACK_TILE 방어: AI가 손패에 없는 타일을 제안할 때 즉시 감지·필터링",
        "• 조커 지원 개선: 조커(JK1/JK2) 포함 조합 제안 능력 향상",
        "• 결과: G1=25.6%(10/39), G2=28.2%(11/39), G3=23.1%(9/39) — 평균 25.6%",
        "• 25초 평균 응답 시간 (CPU 추론, 무료 로컬 실행)",
    ]
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(4.85), Inches(5.7), Inches(1.95))
    tf = tb.text_frame; tf.word_wrap = True
    for j, c in enumerate(changes):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = c; p.font.size = Pt(12); p.font.name = FONT
        p.font.color.rgb = THEME['text']; p.space_after = Pt(5)

    desc_box(slide, [
        "■ 실험 철학",
        "",
        "\"모델이 이해 못 하는",
        "게 아니라 내가 제대로",
        "물어보지 못한 것\"",
        "",
        "이 믿음이 v7, v8, v9를",
        "만들었다.",
        "",
        "■ v10 GO 조건",
        "• ERR_NO_RACK_TILE",
        "  완전 수정",
        "• G1 T73 오류 제거 시",
        "  28%+ 가능",
        "• 조커 확장 지원",
        "",
        "■ 하드웨어",
        "LG Gram 15Z90R",
        "i7-1360P / RAM 16GB",
        "WSL2 + Docker Desktop",
        "",
        "qwen2.5:3b (3B param)",
        "CPU only 추론",
    ], l=6.55, t=0.92, w=3.2, h=6.15)

# ─── 슬라이드 12: UI 아키텍처 ────────────────────────────────────────────────
def slide_ui_arch(prs):
    slide = new_slide(prs)
    slide_header(slide, "UI 아키텍처 — pendingStore SSOT", "Sprint 7 핵심 설계: 38개 파일 영향, 13개 deprecated 필드 제거")

    # 4계층 다이어그램
    layers = [
        ("L4  통신 계층",    "WebSocket / REST API\n서버 ↔ 클라이언트 동기화",       THEME['blue']),
        ("L3  도메인 로직",  "useTurnActions / canConfirmTurn\n순수 함수 — 게임 규칙 적용",  THEME['dark']),
        ("L2  pendingStore", "단일 SSOT — 보드/손패 draft 상태\nhandleDragEnd 9개 분기 통합", THEME['primary']),
        ("L1  순수 UI",      "React 컴포넌트 — 상태 구독\n표시만, 로직 없음",          THEME['secondary']),
    ]
    for i, (label, desc, color) in enumerate(layers):
        y = 1.05 + i * 1.42
        # 계층 박스
        add_round(slide, 0.3, y, 5.9, 1.22, color, line=color)
        add_tb(slide, 0.45, y+0.1, 2.3, 0.45, label, 14, THEME['white'], bold=True)
        add_tb(slide, 0.45, y+0.58, 5.4, 0.58, desc, 12, THEME['white'])
        # 화살표 (아래→위: 이벤트 흐름)
        if i < 3:
            add_arrow(slide, 3.25, y+1.22, 3.25, y+1.42, THEME['accent'], Pt(2))

    # 개선 수치
    metrics = [
        ("38개 파일 수정", THEME['primary']),
        ("13개 필드 제거", THEME['dark']),
        ("9개 분기 통합", THEME['accent']),
        ("단일 SSOT", THEME['blue']),
    ]
    for i, (text, color) in enumerate(metrics):
        x = 0.3 + i * 1.62
        add_round(slide, x, 6.75, 1.52, 0.52, color, line=color)
        add_tb(slide, x, 6.78, 1.52, 0.46, text, 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    desc_box(slide, [
        "■ 이전 아키텍처 문제",
        "",
        "• gameStore + pending",
        "  Store 두 곳에 상태 분산",
        "• 동일 데이터 13개 필드",
        "  중복 관리",
        "• drag 이벤트 처리 로직",
        "  컴포넌트에 혼재",
        "",
        "■ pendingStore 개선",
        "",
        "• draft 단일 SSOT로",
        "  보드+손패 상태 통합",
        "• handleDragEnd 9개",
        "  분기를 한 곳에서 처리",
        "• useTurnActions가",
        "  pendingStore 기반",
        "  순수 함수로 재설계",
        "",
        "■ 핵심 원칙",
        "\"상태는 한 곳에서만",
        "관리한다\"",
        "= pendingStore.draft",
    ], l=6.55, t=0.92, w=3.2, h=6.15)

# ─── 슬라이드 13: 게임룰 SSOT ────────────────────────────────────────────────
def slide_game_rules(prs):
    slide = new_slide(prs)
    slide_header(slide, "게임룰 SSOT — 77개 룰 체계", "55-game-rules-enumeration.md 단일 기준점 — 코드 · 테스트 · UX가 같은 문서를 봄")

    # 룰 카테고리
    categories = [
        ("V-  유효성 검증",   "22개", THEME['primary'],   "그룹/런 조건, 초기 용융\n30점, 조커 규칙, V-22\n단독 승리 룰 등"),
        ("UR- UI 룰",         "40개", THEME['dark'],      "드래그앤드롭 행동 36개\n피드백 색상 규칙\n타이머 16개 지점"),
        ("D-  드로우/패널티", "9개",  THEME['accent'],    "Human=3장, AI=1장\n패널티 B안 확정\n타임아웃 강제 드로우"),
        ("INV-무효화",        "6개",  THEME['blue'],      "잘못된 배치 무효\nERR_NO_RACK_TILE\n재시도 최대 3회"),
    ]
    for i, (cat, count, color, desc) in enumerate(categories):
        x = 0.3 + i * 2.42
        add_round(slide, x, 1.0, 2.2, 4.0, color, line=color)
        add_tb(slide, x+0.08, 1.05, 2.05, 0.45, cat, 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.08, 1.55, 2.05, 0.65, count, 30, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, x+0.2, 2.22, 1.8, 0.04, THEME['white'])
        add_tb(slide, x+0.08, 2.3, 2.05, 2.55, desc, 12, THEME['white'], align=PP_ALIGN.CENTER)

    # 통계
    add_round(slide, 0.3, 5.1, 9.4, 1.85, THEME['light'], line=THEME['primary'], line_w=Pt(1.5))
    stats = [
        ("활성 룰", "77개", THEME['primary']),
        ("결번 룰", "1개 (UR-39 폐기)", THEME['gray']),
        ("행동 매트릭스", "21개 행동", THEME['dark']),
        ("상태 머신", "12개 상태", THEME['dark']),
        ("UX 매핑", "43-rule-ux-sync-ssot", THEME['primary']),
    ]
    for i, (label, val, color) in enumerate(stats):
        x = 0.5 + i * 1.88
        add_tb(slide, x, 5.18, 1.75, 0.38, label, 11, THEME['gray'])
        add_tb(slide, x, 5.55, 1.75, 0.5, val, 13, color, bold=True)
    add_tb(slide, 0.5, 6.15, 9.0, 0.65,
           "게임룰 SSOT가 없으면 구현이 사양을 추월한다. 77개 룰은 코드와 문서가 같은 언어로 말하게 하는 장치다.",
           12, THEME['text'])

# ─── 슬라이드 14: 타임아웃 체인 ──────────────────────────────────────────────
def slide_timeout(prs):
    slide = new_slide(prs)
    slide_header(slide, "타임아웃 체인 — 10개 지점 SSOT", "부등식 계약: script_ws > gs_ctx > http_client > istio_vs > adapter > llm_vendor")

    chain = [
        ("script_ws",    "1070s", "E2E 테스트\n스크립트"),
        ("gs_ctx",       "1060s", "Game Server\nContext"),
        ("http_client",  "1060s", "HTTP 클라이언트\n(Go)"),
        ("istio_vs",     "1010s", "Istio\nVirtualService"),
        ("adapter",      "1000s", "AI Adapter\n내부 한도"),
        ("llm_vendor",   "~900s", "LLM API\n(모델별 상이)"),
    ]
    x_positions = [0.3, 1.95, 3.6, 5.25, 6.9, 8.3]
    for i, ((label, seconds, desc), x) in enumerate(zip(chain, x_positions)):
        color = THEME['primary'] if i < 3 else THEME['dark'] if i < 5 else THEME['accent']
        add_round(slide, x, 1.2, 1.55, 1.3, color, line=color)
        add_tb(slide, x, 1.28, 1.55, 0.45, label, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x, 1.73, 1.55, 0.4, seconds, 16, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x, 2.15, 1.55, 0.3, desc, 9, THEME['white'], align=PP_ALIGN.CENTER)
        if i < 5:
            add_arrow(slide, x+1.55, 1.85, x+1.95, 1.85, THEME['accent'], Pt(2.5))

    # 부등식
    add_round(slide, 0.3, 2.75, 9.4, 0.55, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.5, 2.82, 9.0, 0.42,
           "1070 > 1060 > 1060 > 1010 > 1000 > ~900  →  정상 응답이 fallback으로 오분류되지 않음",
           13, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 사고 사례
    add_round(slide, 0.3, 3.5, 9.4, 2.55, THEME['light'], line=THEME['red'], line_w=Pt(2))
    add_tb(slide, 0.5, 3.55, 9.0, 0.4, "2026-04-16 사고: 타임아웃 체인 부등식 위반", 14, THEME['red'], bold=True)
    incident = [
        "원인: adapter 내부 한도(700s) < istio_vs(1010s) — 계약 위반",
        "증상: LLM이 정상 응답을 반환했지만 adapter가 먼저 타임아웃 → fallback으로 오분류",
        "조치: adapter 1000s로 상향, 타임아웃 레지스트리 문서(41번) SSOT 확립",
        "교훈: 한 값을 바꿀 때 반드시 10개 지점 전수 체크 (41번 문서 §3+§5)",
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.98), Inches(9.0), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(incident):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "• " + line; p.font.size = Pt(12); p.font.name = FONT
        p.font.color.rgb = THEME['text']; p.space_after = Pt(5)

    add_round(slide, 0.3, 6.15, 9.4, 0.75, THEME['primary'], line=THEME['primary'])
    add_tb(slide, 0.5, 6.22, 9.0, 0.58,
           "SSOT: docs/02-design/41-timeout-chain-breakdown.md  |  ConfigMap: AI_ADAPTER_TIMEOUT_SEC=1000",
           13, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── 슬라이드 15: 품질 지표 ──────────────────────────────────────────────────
def slide_quality(prs):
    slide = new_slide(prs)
    slide_header(slide, "품질 지표 — 테스트 2,462건 최종 현황")

    # 테스트 현황 카드
    test_data = [
        ("Frontend\nJest",      "659",  "PASS",  THEME['primary']),
        ("Go 단위\n테스트",      "770",  "PASS",  THEME['dark']),
        ("AI Adapter\n(NestJS)", "637",  "PASS",  THEME['primary']),
        ("E2E\n(Playwright)",    "396",  "PASS",  THEME['blue']),
    ]
    for i, (label, count, status, color) in enumerate(test_data):
        x = 0.3 + i * 2.42
        add_round(slide, x, 1.0, 2.2, 2.5, color, line=color)
        add_tb(slide, x+0.08, 1.08, 2.05, 0.55, label, 13, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.08, 1.65, 2.05, 0.9, count, 36, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.08, 2.6, 2.05, 0.42, status, 14, THEME['white'], align=PP_ALIGN.CENTER)

    # 총계
    add_round(slide, 0.3, 3.65, 9.4, 0.7, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.5, 3.72, 9.0, 0.55,
           "총계: 2,462 PASS  /  0 FAIL  — 63일 전체 기간 FAIL 0건 유지",
           18, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # CI/CD
    add_round(slide, 0.3, 4.5, 4.55, 2.4, THEME['light'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.5, 4.55, 4.2, 0.4, "CI/CD 파이프라인", 14, THEME['dark'], bold=True)
    ci_lines = [
        "• GitLab CI 17 stages — ALL GREEN (Pipeline #113+)",
        "• Trivy: Critical/High CVE = 0건 유지",
        "• SonarQube: 3개 프로젝트 Quality Gate PASS",
        "• Kaniko 빌드 → ArgoCD GitOps 자동 배포",
        "• K8s 7개 서비스 Running (rummikub ns)",
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(4.2), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(ci_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12); p.font.name = FONT
        p.font.color.rgb = THEME['text']; p.space_after = Pt(4)

    # 보안
    add_round(slide, 5.15, 4.5, 4.55, 2.4, THEME['light'], line=THEME['dark'], line_w=Pt(1.5))
    add_tb(slide, 5.35, 4.55, 4.2, 0.4, "보안 성과", 14, THEME['dark'], bold=True)
    sec_lines = [
        "• SEC-A: Go 1.25.9 + go-redis v9.7.3 (vuln 25→0)",
        "• SEC-B: Rate Limiting 강화",
        "• SEC-C: 인증/인가 완전 분리 (OAuth=identity only)",
        "• JWT fail-fast, dev-login 차단",
        "• Critical/High CVE 63일 내내 0건",
    ]
    tb2 = slide.shapes.add_textbox(Inches(5.35), Inches(5.0), Inches(4.2), Inches(1.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for j, line in enumerate(sec_lines):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line; p.font.size = Pt(12); p.font.name = FONT
        p.font.color.rgb = THEME['text']; p.space_after = Pt(4)

# ─── 슬라이드 16: 운영 환경 ──────────────────────────────────────────────────
def slide_k8s(prs):
    slide = new_slide(prs)
    slide_header(slide, "K8s 운영 환경 — 현재 배포 상태", "namespace: rummikub | 7개 서비스 Running")

    services = [
        ("frontend",    "rummikub-frontend:lobby-fix-e7222d0",          "Next.js 게임 UI",      THEME['primary']),
        ("game-server", "rummiarena/game-server:day5-8dc0999",           "Go 게임 엔진 + WS",   THEME['dark']),
        ("ai-adapter",  "rummiarena/ai-adapter:v9-ollama-place-8631831", "NestJS LLM 연동",      THEME['primary']),
        ("ollama",      "ollama/ollama:latest (qwen2.5:3b PVC)",         "로컬 LLM",             THEME['accent']),
        ("postgres",    "postgres:16",                                    "영속 데이터 DB",       THEME['blue']),
        ("redis",       "redis:7-alpine",                                 "게임 상태 캐시",       THEME['blue']),
        ("admin",       "rummiarena/admin:latest",                        "관리자 대시보드",      THEME['gray']),
    ]
    for i, (name, image, desc, color) in enumerate(services):
        y = 1.05 + i * 0.87
        add_round(slide, 0.3, y, 1.5, 0.72, color, line=color)
        add_tb(slide, 0.3, y+0.14, 1.5, 0.45, name, 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_round(slide, 1.9, y+0.08, 4.55, 0.56, THEME['light'], line=THEME['gray'], line_w=Pt(1))
        add_tb(slide, 2.05, y+0.12, 4.35, 0.44, image, 9.5, THEME['text'])
        add_round(slide, 6.55, y+0.08, 2.5, 0.56, THEME['light'], line=color, line_w=Pt(1))
        add_tb(slide, 6.65, y+0.12, 2.35, 0.44, desc, 11, THEME['text'])
        # RUNNING 배지
        s = add_round(slide, 9.15, y+0.14, 0.62, 0.44, THEME['primary'], line=THEME['primary'])
        add_tb(slide, 9.15, y+0.18, 0.62, 0.36, "RUN", 9, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 설정
    add_round(slide, 0.3, 7.17, 9.4, 0.4, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.5, 7.22, 9.0, 0.32,
           "ConfigMap: AI_ADAPTER_TIMEOUT_SEC=1000  |  GAME_MAX_TURNS_LIMIT=200  |  Istio VS ai-adapter timeout=1010s",
           11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── 슬라이드 17: 동시접속 분석 ──────────────────────────────────────────────
def slide_concurrent(prs):
    slide = new_slide(prs)
    slide_header(slide, "동시접속 분석 — 100~200명 시나리오", "단일 노드(LG Gram) 기준 현실적 한계와 클라우드 확장 방안")

    scenarios = [
        ("A", "Human Only",     "200명 가능",    THEME['primary'],
         "• game-server replicas=3\n• 방당 2~4명 × 50~100방\n• Redis/Postgres 충분\n• CPU 병목 없음"),
        ("B", "API LLM 혼합",   "10게임 동시",   THEME['dark'],
         "• DeepSeek/GPT 100초 블로킹\n• 비용 제약: ~$0.4/게임\n• 10게임 × 4턴 = 동시 처리\n• 40 LLM 호출/시"),
        ("C", "Ollama 로컬",     "2~3게임 한계",  THEME['accent'],
         "• CPU 단일 스레드 추론\n• qwen2.5:3b 25초/턴\n• 동시 처리 불가 (큐잉)\n• GPU 있으면 10x 향상"),
        ("D", "클라우드(GKE)",   "$180~280/월",   THEME['blue'],
         "• 3× game-server pods\n• GPU node (Ollama)\n• 100~200명 안정 지원\n• 예상 비용 $180~280/월"),
    ]
    for i, (alpha, scenario, result, color, detail) in enumerate(scenarios):
        col = i % 2; row = i // 2
        x = 0.3 + col * 4.85; y = 1.0 + row * 3.2
        add_round(slide, x, y, 4.6, 3.0, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, x, y, 4.6, 0.58, color)
        add_tb(slide, x+0.1, y+0.08, 0.5, 0.44, alpha, 18, THEME['white'], bold=True)
        add_tb(slide, x+0.65, y+0.08, 2.6, 0.44, scenario, 14, THEME['white'], bold=True)
        add_tb(slide, x+3.3, y+0.12, 1.2, 0.36, result, 13, THEME['white'], bold=True, align=PP_ALIGN.RIGHT)
        tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.68), Inches(4.2), Inches(2.2))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for line in detail.split('\n'):
            if first: p = tf.paragraphs[0]; first = False
            else: p = tf.add_paragraph()
            p.text = line; p.font.size = Pt(13); p.font.name = FONT
            p.font.color.rgb = THEME['text']; p.space_after = Pt(6)

# ─── 슬라이드 18: 기술 부채 ──────────────────────────────────────────────────
def slide_tech_debt(prs):
    slide = new_slide(prs)
    slide_header(slide, "기술 부채 — 미해결 항목 및 향후 과제", "프로젝트 종료 시점 식별된 부채 — 다음 버전 착수 시 우선 처리")

    debts = [
        ("SEC-DEBT-001", "HIGH",   "next-auth v4→v5 이주 미완",               "frontend-dev",   THEME['red']),
        ("SEC-DEBT-003", "HIGH",   "WebSocket 인증 토큰 갱신 미구현\n(만료된 토큰 사용 가능성)", "go-dev", THEME['red']),
        ("SEC-DEBT-004", "MED",    "Admin API Rate Limiting 미적용",            "node-dev",       THEME['accent']),
        ("SEC-DEBT-005", "MED",    "Refresh Token Rotation 미구현",             "go-dev",         THEME['accent']),
        ("P3-3",         "LOW",    "DndContext GameRoom 이전 미완",             "frontend-dev",   THEME['primary']),
        ("v10 Ollama",   "FUTURE", "ERR_NO_RACK_TILE 수정 → 30%+ 목표",        "ai-engineer",    THEME['dark']),
    ]
    for i, (id_, priority, desc, owner, color) in enumerate(debts):
        y = 1.05 + i * 1.02
        # 우선순위 배지
        add_round(slide, 0.3, y+0.08, 0.9, 0.56, color, line=color)
        add_tb(slide, 0.3, y+0.12, 0.9, 0.48, priority, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        # ID
        add_round(slide, 1.3, y+0.06, 1.4, 0.6, THEME['light'], line=color, line_w=Pt(1.5))
        add_tb(slide, 1.35, y+0.12, 1.3, 0.48, id_, 12, THEME['text'], bold=True)
        # 설명
        add_tb(slide, 2.85, y+0.06, 4.85, 0.72, desc, 12, THEME['text'])
        # 담당
        add_round(slide, 7.8, y+0.12, 1.8, 0.48, THEME['light'], line=THEME['gray'], line_w=Pt(1))
        add_tb(slide, 7.85, y+0.16, 1.7, 0.38, owner, 11, THEME['gray'])

    add_round(slide, 0.3, 7.17, 9.4, 0.4, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.5, 7.22, 9.0, 0.32,
           "식별된 부채는 반드시 타임라인과 처리 계획이 있어야 한다 — '나중에'는 없다",
           13, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── 슬라이드 19: 최종 성과 ──────────────────────────────────────────────────
def slide_achievements(prs):
    slide = new_slide(prs)
    slide_header(slide, "프로젝트 최종 성과 — 63일의 결산")

    # 큰 숫자들
    big_stats = [
        ("750+",   "커밋",          THEME['primary']),
        ("110+",   "문서",          THEME['dark']),
        ("2,462",  "테스트 PASS",   THEME['accent']),
        ("77",     "게임룰 SSOT",   THEME['blue']),
        ("33.3%",  "AI 최고 배치율", THEME['primary']),
        ("17/17",  "CI/CD 스테이지", THEME['dark']),
    ]
    for i, (num, label, color) in enumerate(big_stats):
        col = i % 3; row = i // 3
        x = 0.35 + col * 3.15; y = 1.0 + row * 2.2
        add_round(slide, x, y, 2.9, 1.95, color, line=color)
        add_tb(slide, x+0.1, y+0.2, 2.7, 0.85, num, 34, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.1, y+1.1, 2.7, 0.7, label, 14, THEME['white'], align=PP_ALIGN.CENTER)

    # 요약
    add_round(slide, 0.3, 5.5, 9.4, 1.65, THEME['light'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.5, 5.55, 9.0, 0.4, "이 프로젝트가 증명한 것", 14, THEME['dark'], bold=True)
    proof_lines = [
        "• 13명의 에이전트가 63일 동안 단일 프로젝트를 완성할 수 있다",
        "• 3B 파라미터 로컬 LLM도 프롬프트 엔지니어링으로 25.6% 배치율을 달성할 수 있다",
        "• 아키텍처 설계가 선행될 때 Phase D처럼 8.5 Story Point를 1.5시간에 처리할 수 있다",
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.98), Inches(9.0), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(proof_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12); p.font.name = FONT
        p.font.color.rgb = THEME['text']; p.space_after = Pt(3)

# ─── 슬라이드 20: 마무리 ─────────────────────────────────────────────────────
def slide_closing(prs):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, 10, 7.5, THEME['dark'])
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(4.5), Inches(5), Inches(4))
    s.fill.solid(); s.fill.fore_color.rgb = THEME['primary']; s.line.color.rgb = THEME['primary']

    add_tb(slide, 1.0, 1.0, 8.0, 0.8,
           "RummiArena 프로젝트 종료", 36, THEME['secondary'], bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, 1.85, 7.0, 0.06, THEME['secondary'])
    add_tb(slide, 1.0, 2.0, 8.0, 0.6,
           "2026-03-08 ~ 2026-05-10  (63일)", 20, THEME['light'], align=PP_ALIGN.CENTER)

    # 팀에게
    add_tb(slide, 1.0, 2.8, 8.0, 1.4,
           "13명의 에이전트가 함께 달렸습니다.\narchinect의 설계력, go-dev의 안정성, frontend-dev의 실행력,\nqa의 의심, ai-engineer의 끈기, 그리고 모든 팀원의 헌신이\n750+ 커밋과 2,462 테스트를 만들었습니다.",
           14, THEME['white'], align=PP_ALIGN.CENTER)

    add_rect(slide, 1.5, 4.3, 7.0, 0.06, THEME['secondary'])

    add_tb(slide, 1.0, 4.5, 8.0, 0.6,
           "\"계획은 실행 시작 전까지만 유효하다\"", 16, THEME['secondary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 1.0, 5.1, 8.0, 0.55,
           "\"Merge Gate는 속도를 늦추지 않는다\"", 16, THEME['secondary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 1.0, 5.65, 8.0, 0.55,
           "\"데일리 로그와 세션 로그가 기억을 대체한다\"", 16, THEME['secondary'], align=PP_ALIGN.CENTER)

    add_tb(slide, 1.0, 6.4, 8.0, 0.5,
           "— pm 최종 회고 中 —  RummiArena 2026-05-10", 13, THEME['light'], align=PP_ALIGN.CENTER)

# ─── 메인 ────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("슬라이드 생성 중...")
    slide_title(prs);         print("  1/20 표지")
    slide_toc(prs);           print("  2/20 목차")
    slide_overview(prs);      print("  3/20 프로젝트 개요")
    slide_team(prs);          print("  4/20 팀 구성")
    slide_architecture(prs);  print("  5/20 시스템 아키텍처")
    slide_tech_stack(prs);    print("  6/20 기술 스택")
    slide_sprint_timeline(prs); print("  7/20 스프린트 타임라인")
    slide_sprint_1_4(prs);    print("  8/20 Sprint 1~4")
    slide_sprint_5_7(prs);    print("  9/20 Sprint 5~7 + 핫픽스")
    slide_llm_results(prs);   print(" 10/20 AI 대전 결과")
    slide_ollama(prs);        print(" 11/20 Ollama 진화")
    slide_ui_arch(prs);       print(" 12/20 UI 아키텍처")
    slide_game_rules(prs);    print(" 13/20 게임룰 SSOT")
    slide_timeout(prs);       print(" 14/20 타임아웃 체인")
    slide_quality(prs);       print(" 15/20 품질 지표")
    slide_k8s(prs);           print(" 16/20 K8s 운영 환경")
    slide_concurrent(prs);    print(" 17/20 동시접속 분석")
    slide_tech_debt(prs);     print(" 18/20 기술 부채")
    slide_achievements(prs);  print(" 19/20 최종 성과")
    slide_closing(prs);       print(" 20/20 마무리")

    out = os.path.join(os.path.dirname(__file__), "RummiArena_프로젝트종료보고서.pptx")
    prs.save(out)
    print(f"\n완료: {out}")
    print(f"총 슬라이드: {len(prs.slides)}장")

if __name__ == "__main__":
    main()
