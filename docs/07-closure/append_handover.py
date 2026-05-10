"""
RummiArena 운영 이관 계획 — 별첨 슬라이드 추가
기존 RummiArena_프로젝트종료보고서.pptx 에 append
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os

# ─── 테마 (동일) ─────────────────────────────────────────────────────────────
THEME = {
    'primary':   RGBColor(46, 125, 50),
    'secondary': RGBColor(129, 199, 132),
    'accent':    RGBColor(255, 152, 0),
    'dark':      RGBColor(27, 94, 32),
    'light':     RGBColor(232, 245, 233),
    'text':      RGBColor(33, 33, 33),
    'white':     RGBColor(255, 255, 255),
    'gray':      RGBColor(120, 120, 120),
    'lgray':     RGBColor(200, 200, 200),
    'red':       RGBColor(198, 40, 40),
    'lred':      RGBColor(255, 235, 235),
    'blue':      RGBColor(21, 101, 192),
    'lblue':     RGBColor(227, 242, 253),
    'yellow':    RGBColor(245, 127, 23),
    'lyellow':   RGBColor(255, 248, 225),
    'purple':    RGBColor(106, 27, 154),
    'desc_bg':   RGBColor(226, 240, 217),
}
FONT = "맑은 고딕"

# ─── 헬퍼 ────────────────────────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill; s.line.width = line_w
    return s

def add_round(slide, l, t, w, h, fill, line=None, line_w=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill; s.line.width = line_w
    return s

def add_diam(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill; s.line.width = Pt(1.5)
    return s

def add_para(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill; s.line.width = Pt(1.5)
    return s

def add_hex(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.HEXAGON,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill; s.line.width = Pt(1.5)
    return s

def add_tb(slide, l, t, w, h, text, size, color,
           bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FONT
    p.font.italic = italic; p.alignment = align
    return tb

def add_tb_lines(slide, l, t, w, h, lines, size, color,
                 bold=False, space=Pt(4), indent_char=""):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = indent_char + line
        p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = FONT; p.space_after = space
    return tb

def add_arrow(slide, x1, y1, x2, y2, color=None, w=Pt(2.5), arrow=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color or THEME['accent']; c.line.width = w
    if arrow: c.line.end_arrow_type = 2
    return c

def hdr(slide, title, subtitle=None, color=None):
    add_rect(slide, 0, 0, 10, 0.85, color or THEME['dark'])
    add_tb(slide, 0.35, 0.1, 7.5, 0.65, title, 24, THEME['white'], bold=True)
    if subtitle:
        add_tb(slide, 0.35, 0.57, 9.3, 0.3, subtitle, 11, THEME['secondary'])
    add_rect(slide, 0, 7.22, 10, 0.05, THEME['secondary'])
    add_tb(slide, 0.3, 7.27, 9.4, 0.2,
           "【별첨】RummiArena 운영 이관 계획  |  2026-05-10",
           9, THEME['gray'])
    # 별첨 배지
    s = add_round(slide, 8.85, 0.12, 0.88, 0.56, THEME['accent'], line=THEME['accent'])
    add_tb(slide, 8.85, 0.15, 0.88, 0.5, "별 첨", 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

def desc_box(slide, lines, l=6.6, t=0.93, w=3.2, h=6.1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = THEME['desc_bg']
    s.line.color.rgb = THEME['desc_bg']
    tb = slide.shapes.add_textbox(Inches(l+0.12), Inches(t+0.12),
                                   Inches(w-0.24), Inches(h-0.24))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("■"):
            p.text = line; p.font.size = Pt(10); p.font.bold = True
            p.font.color.rgb = THEME['dark']
        elif line.startswith("※"):
            p.text = line; p.font.size = Pt(9); p.font.italic = True
            p.font.color.rgb = THEME['gray']
        else:
            p.text = line; p.font.size = Pt(9.5)
            p.font.color.rgb = THEME['text']
        p.font.name = FONT; p.space_after = Pt(2.5)

# ══════════════════════════════════════════════════════════════════════════════
# 별첨 슬라이드 정의
# ══════════════════════════════════════════════════════════════════════════════

# ─── A-0: 별첨 표지 ──────────────────────────────────────────────────────────
def slide_annex_cover(prs):
    slide = new_slide(prs)
    add_rect(slide, 0, 0, 10, 7.5, THEME['dark'])
    # 대각선 장식
    for i in range(5):
        s = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(7.5 + i*0.6), Inches(0), Inches(0.5), Inches(7.5))
        s.fill.solid(); s.fill.fore_color.rgb = THEME['primary']
        s.line.color.rgb = THEME['primary']
    add_rect(slide, 0.5, 2.8, 9, 0.08, THEME['accent'])
    add_tb(slide, 0.8, 1.1, 8.5, 0.7,
           "【별첨】", 22, THEME['secondary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 0.8, 1.8, 8.5, 1.1,
           "RummiArena 운영 이관 계획", 42, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 0.8, 3.0, 8.5, 0.6,
           "Operations Handover Plan", 22, THEME['secondary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 0.8, 3.75, 8.5, 0.5,
           "현재 운영 환경 현황 · 이관 절차 · 장애 대응 · 스케일링 방안",
           16, THEME['light'], align=PP_ALIGN.CENTER)
    sections = ["A-1  현재 운영 환경 아키텍처",
                "A-2  K8s 서비스 구조도",
                "A-3  GitOps 배포 파이프라인",
                "A-4  서비스 시작 순서 (의존성)",
                "A-5  시크릿 관리 구조",
                "A-6  일상 운영 절차",
                "A-7  장애 대응 플로우",
                "A-8  장애 유형별 대응 (1/2)",
                "A-9  장애 유형별 대응 (2/2)",
                "A-10 배포 절차 (GitOps/수동/롤백)",
                "A-11 스케일링 전략",
                "A-12 클라우드 이관 아키텍처 (GKE)",
                "A-13 이관 전 체크리스트",
                "A-14 이관 후 검증 절차"]
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, s in enumerate(sections):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s; p.font.size = Pt(11); p.font.color.rgb = THEME['light']
        p.font.name = FONT; p.space_after = Pt(1)

# ─── A-1: 현재 운영 환경 전체 아키텍처 ──────────────────────────────────────
def slide_env_arch(prs):
    slide = new_slide(prs)
    hdr(slide, "A-1  현재 운영 환경 아키텍처",
        "Windows 11 + WSL2 + Docker Desktop K8s (단일 노드)")

    # 최상단: 호스트 OS 레이어
    add_round(slide, 0.25, 0.95, 9.5, 6.2, THEME['lgray'], line=THEME['gray'], line_w=Pt(1))
    add_tb(slide, 0.35, 0.97, 3.0, 0.32, "Windows 11 Host (LG Gram 15Z90R · i7-1360P · RAM 16GB · SSD)",
           9, THEME['gray'], italic=True)

    # WSL2 레이어
    add_round(slide, 0.45, 1.3, 9.1, 5.7, THEME['light'], line=THEME['secondary'], line_w=Pt(1.5))
    add_tb(slide, 0.55, 1.32, 3.0, 0.28, "WSL2 (Ubuntu 22.04 · 메모리 10GB)",
           9, THEME['primary'], italic=True)

    # Docker Desktop
    add_round(slide, 0.65, 1.65, 8.7, 5.1, THEME['white'], line=THEME['primary'], line_w=Pt(2))
    add_tb(slide, 0.75, 1.67, 4.0, 0.28, "Docker Desktop K8s  |  namespace: rummikub",
           9.5, THEME['dark'], bold=True, italic=True)

    # 외부 사용자 (왼쪽)
    add_round(slide, 0.15, 2.3, 1.2, 0.65, THEME['secondary'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.15, 2.35, 1.2, 0.55, "사용자\nBrowser", 10, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 0.15, 3.15, 1.2, 0.65, THEME['secondary'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.15, 3.2, 1.2, 0.55, "관리자\nBrowser", 10, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)

    # Traefik Ingress
    add_rect(slide, 1.55, 2.2, 1.4, 1.4, THEME['primary'])
    add_tb(slide, 1.55, 2.28, 1.4, 0.5, "Traefik\nIngress", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 1.55, 2.78, 1.4, 0.7, ":30000\n:30001\n:30080", 9, THEME['light'], align=PP_ALIGN.CENTER)
    add_arrow(slide, 1.35, 2.63, 1.55, 2.63, THEME['primary'])
    add_arrow(slide, 1.35, 3.48, 1.55, 3.2, THEME['primary'])

    # Frontend / Admin
    add_round(slide, 3.2, 2.05, 1.5, 0.75, THEME['primary'])
    add_tb(slide, 3.2, 2.1, 1.5, 0.65, "Frontend\nNext.js :3000", 10, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 3.2, 2.95, 1.5, 0.75, THEME['dark'])
    add_tb(slide, 3.2, 3.0, 1.5, 0.65, "Admin\nNext.js :3001", 10, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 2.95, 2.42, 3.2, 2.42, THEME['accent'])
    add_arrow(slide, 2.95, 3.33, 3.2, 3.33, THEME['accent'])

    # Game Server (중앙)
    add_round(slide, 5.0, 1.95, 1.7, 1.5, THEME['dark'])
    add_tb(slide, 5.0, 2.02, 1.7, 0.42, "Game Server", 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 5.0, 2.42, 1.7, 0.95,
           "Go / gin\ngorilla/ws\n:8080", 10, THEME['light'], align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.7, 2.42, 5.0, 2.42, THEME['accent'])
    add_arrow(slide, 4.7, 3.33, 5.0, 3.22, THEME['accent'])

    # AI Adapter
    add_round(slide, 7.0, 2.1, 1.55, 1.3, THEME['primary'])
    add_tb(slide, 7.0, 2.17, 1.55, 0.4, "AI Adapter", 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, 7.0, 2.55, 1.55, 0.75, "NestJS\n:3001", 10, THEME['white'], align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.7, 2.7, 7.0, 2.7, THEME['accent'])

    # LLM 클라우드
    llms = [("OpenAI\nGPT-5-mini", 1.95), ("DeepSeek\nV4-Pro", 2.7), ("Claude\nSonnet 4", 3.45)]
    for label, yy in llms:
        add_round(slide, 8.8, yy, 1.15, 0.6, THEME['accent'], line=THEME['accent'])
        add_tb(slide, 8.8, yy+0.04, 1.15, 0.52, label, 9, THEME['white'], align=PP_ALIGN.CENTER)
        add_arrow(slide, 8.55, yy+0.3, 8.8, yy+0.3, THEME['secondary'], Pt(1.5))

    # Ollama (로컬 K8s)
    add_round(slide, 7.15, 3.7, 2.4, 0.75, THEME['yellow'], line=THEME['accent'])
    add_tb(slide, 7.15, 3.75, 2.4, 0.65, "Ollama (K8s Pod)\nqwen2.5:3b  CPU", 10, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 7.62, 3.4, 7.62, 3.7, THEME['accent'])

    # 데이터 레이어
    add_round(slide, 1.0, 4.75, 2.2, 0.82, THEME['blue'])
    add_tb(slide, 1.0, 4.8, 2.2, 0.72, "Redis 7\n게임 상태 · 세션\n:6379", 9.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 3.6, 4.75, 2.2, 0.82, THEME['blue'])
    add_tb(slide, 3.6, 4.8, 2.2, 0.72, "PostgreSQL 16\n영속 데이터 (ELO 등)\n:5432", 9.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 6.2, 4.75, 2.0, 0.82, THEME['gray'])
    add_tb(slide, 6.2, 4.8, 2.0, 0.72, "Istio Sidecar\ngame-server\nai-adapter", 9.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 8.45, 4.75, 1.25, 0.82, THEME['gray'])
    add_tb(slide, 8.45, 4.8, 1.25, 0.72, "ArgoCD\nGitOps\n:8080", 9.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 연결선 (Game Server → DB)
    add_arrow(slide, 5.85, 3.45, 5.85, 4.75, THEME['blue'], Pt(1.5))
    add_arrow(slide, 5.65, 3.45, 2.1, 4.75, THEME['blue'], Pt(1.5))

    # NodePort 라벨
    add_rect(slide, 0.65, 6.0, 8.7, 0.35, THEME['dark'])
    add_tb(slide, 0.75, 6.03, 8.5, 0.3,
           "NodePort: 30000(FE) · 30001(Admin) · 30080(GS) · 30081(AI) · 30432(PG)  |  Helm 차트 5개",
           11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── A-2: K8s 서비스 구조도 ──────────────────────────────────────────────────
def slide_k8s_detail(prs):
    slide = new_slide(prs)
    hdr(slide, "A-2  K8s 서비스 구조도",
        "namespace: rummikub  |  Deployment · Service · PVC · ConfigMap · Secret")

    # 컬럼 헤더
    cols = ["Deployment", "Service (NodePort)", "Volume/Secret", "이미지 태그"]
    col_x = [0.28, 2.85, 5.55, 7.5]
    col_w = [2.5, 2.6, 1.85, 2.35]
    for i, (label, x, w) in enumerate(zip(cols, col_x, col_w)):
        add_round(slide, x, 0.93, w, 0.42, THEME['primary'], line=THEME['primary'])
        add_tb(slide, x+0.05, 0.95, w-0.1, 0.38, label, 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    rows = [
        ("frontend",     "ClusterIP :3000\nNodePort 30000",  "Secret: nextauth\ngoogle-oauth",     "lobby-fix-e7222d0",        THEME['primary']),
        ("admin",        "ClusterIP :3001\nNodePort 30001",  "Secret: nextauth\ngoogle-oauth",     "latest",                   THEME['dark']),
        ("game-server",  "ClusterIP :8080\nNodePort 30080",  "ConfigMap: main\nSecret: db, llm",   "day5-8dc0999",             THEME['dark']),
        ("ai-adapter",   "ClusterIP :3001\nNodePort 30081",  "ConfigMap: main\nSecret: llm-api",   "v9-ollama-place-8631831",  THEME['primary']),
        ("postgres",     "ClusterIP :5432\nNodePort 30432",  "PVC: 5Gi\nSecret: db-secret",        "postgres:16-alpine",       THEME['blue']),
        ("redis",        "ClusterIP :6379\n(내부 전용)",      "PVC: 1Gi\n(선택)",                   "redis:7-alpine",           THEME['blue']),
        ("ollama",       "ClusterIP :11434\n(내부 전용)",     "PVC: 10Gi\n(모델 파일)",              "ollama/ollama:latest",     RGBColor(180, 100, 0)),
    ]
    for i, (dep, svc, vol, img, color) in enumerate(rows):
        y = 1.45 + i * 0.84
        # 배포명
        add_round(slide, col_x[0], y+0.05, col_w[0], 0.72, color, line=color)
        add_tb(slide, col_x[0]+0.1, y+0.1, col_w[0]-0.2, 0.62, dep, 13, THEME['white'], bold=True)
        # 서비스
        add_round(slide, col_x[1], y+0.05, col_w[1], 0.72, THEME['light'], line=color, line_w=Pt(1.5))
        add_tb(slide, col_x[1]+0.1, y+0.1, col_w[1]-0.2, 0.62, svc, 11, THEME['text'])
        # 볼륨/시크릿
        add_round(slide, col_x[2], y+0.05, col_w[2], 0.72, THEME['light'], line=THEME['gray'], line_w=Pt(1))
        add_tb(slide, col_x[2]+0.08, y+0.1, col_w[2]-0.16, 0.62, vol, 10, THEME['text'])
        # 이미지 태그
        add_round(slide, col_x[3], y+0.05, col_w[3], 0.72, THEME['light'], line=THEME['gray'], line_w=Pt(1))
        add_tb(slide, col_x[3]+0.08, y+0.1, col_w[3]-0.16, 0.62, img, 9.5, THEME['blue'])
        # Running 배지
        s = add_round(slide, 9.87, y+0.18, 0.08, 0.44, THEME['primary'])  # 녹색 점
        add_round(slide, 9.82, y+0.18, 0.15, 0.44, THEME['primary'], line=THEME['primary'])

    # 요약
    add_round(slide, 0.25, 7.18, 9.55, 0.4, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.45, 7.22, 9.2, 0.32,
           "ConfigMap: AI_ADAPTER_TIMEOUT_SEC=1000  GAME_MAX_TURNS_LIMIT=200  OLLAMA_PROMPT_VARIANT=v9-ollama-place  USE_V2_PROMPT=true",
           10.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── A-3: GitOps 배포 파이프라인 ─────────────────────────────────────────────
def slide_gitops(prs):
    slide = new_slide(prs)
    hdr(slide, "A-3  GitOps 배포 파이프라인",
        "코드 변경 → GitHub(SSOT) → GitLab(Mirror) → CI 17단계 → ArgoCD → K8s")

    # ── 레인 배경 ──
    lanes = [
        ("개발자 (애벌레)",   THEME['light'],    0.93, 1.18),
        ("GitHub (SSOT)",   THEME['lblue'],    2.15, 1.18),
        ("GitLab CI",       THEME['lyellow'],  3.35, 2.35),
        ("ArgoCD",          THEME['light'],    5.7,  1.0),
        ("K8s Cluster",     THEME['light'],    6.75, 0.85),
    ]
    lane_colors = [THEME['light'], THEME['lblue'], THEME['lyellow'], THEME['light'], THEME['light']]
    lane_y = [0.93, 2.15, 3.35, 5.7, 6.75]
    lane_h = [1.18, 1.18, 2.35, 1.0, 0.85]
    lane_labels = ["개발자 (애벌레)", "GitHub (SSOT)", "GitLab CI (17 Stages)", "ArgoCD", "K8s Cluster"]
    lane_border = [THEME['primary'], THEME['blue'], THEME['yellow'], THEME['dark'], THEME['primary']]

    for label, bg, y, h, border in zip(lane_labels, lane_colors, lane_y, lane_h, lane_border):
        add_round(slide, 0.2, y, 9.55, h, bg, line=border, line_w=Pt(1))
        add_tb(slide, 0.28, y+0.06, 1.2, 0.3, label, 9, border, bold=True, italic=True)

    # ── 박스들 ──
    # 개발자
    add_round(slide, 1.5, 1.02, 1.4, 0.8, THEME['primary'])
    add_tb(slide, 1.5, 1.12, 1.4, 0.6, "git push\norigin main", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # GitHub
    add_round(slide, 3.1, 2.24, 1.4, 0.8, THEME['blue'])
    add_tb(slide, 3.1, 2.34, 1.4, 0.6, "GitHub\nSSO Pull", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, 5.0, 2.24, 1.8, 0.8, THEME['blue'])
    add_tb(slide, 5.0, 2.34, 1.8, 0.6, "GitLab Mirror\n(자동 동기화)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # CI 단계들 (17 stages 압축 표현)
    ci_stages = [
        ("lint\ntest", 1.5, 3.55),
        ("build\n(Kaniko)", 3.0, 3.55),
        ("scan\nTrivy+SQ", 4.5, 3.55),
        ("deploy\nHelm", 6.0, 3.55),
        ("verify\nE2E", 7.5, 3.55),
    ]
    for label, x, y in ci_stages:
        add_round(slide, x, y, 1.25, 0.78, THEME['yellow'], line=RGBColor(180, 100, 0))
        add_tb(slide, x, y+0.1, 1.25, 0.58, label, 11, THEME['dark'], bold=True, align=PP_ALIGN.CENTER)
    # CI 화살표
    for x in [2.75, 4.25, 5.75, 7.25]:
        add_arrow(slide, x, 3.94, x+0.25, 3.94, RGBColor(180, 100, 0), Pt(2))
    # GitLab → CI
    add_arrow(slide, 5.9, 3.04, 5.9, 3.55, THEME['yellow'], Pt(2))
    # GitHub trigger
    add_arrow(slide, 2.9, 2.63, 3.1, 2.63, THEME['blue'], Pt(2))
    add_arrow(slide, 4.5, 2.63, 5.0, 2.63, THEME['blue'], Pt(2))
    # 개발자 → GitHub
    add_arrow(slide, 2.9, 1.42, 3.1, 2.42, THEME['primary'], Pt(2))

    # ArgoCD
    add_round(slide, 3.5, 5.8, 2.0, 0.7, THEME['dark'])
    add_tb(slide, 3.5, 5.9, 2.0, 0.5, "ArgoCD\n(Helm sync)", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 8.75, 4.33, 8.75, 5.3, THEME['dark'], Pt(2))
    add_arrow(slide, 8.75, 5.3, 5.5, 6.15, THEME['dark'], Pt(2))

    # K8s
    add_round(slide, 2.0, 6.83, 5.5, 0.55, THEME['primary'])
    add_tb(slide, 2.0, 6.88, 5.5, 0.45,
           "K8s Deployments  →  롤링 업데이트  →  헬스체크 통과  →  완료", 12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.5, 6.5, 4.5, 6.83, THEME['primary'], Pt(2))

    desc_box(slide, [
        "■ CI 17단계 요약",
        "",
        "1. Static Analysis (lint)",
        "2. Unit Test (Go/Jest)",
        "3. Integration Test",
        "4. Build (Kaniko)",
        "5. Trivy 보안 스캔",
        "6. SonarQube 품질",
        "7. Helm template 검증",
        "8. K8s 배포",
        "9~17. E2E, 성능, 알림",
        "",
        "■ GitOps 원칙",
        "• GitHub = 단일 SSOT",
        "• GitLab = CI 실행 미러",
        "• ArgoCD = 배포 게이트",
        "• 수동 kubectl apply",
        "  = 금지 (긴급 시 예외)",
        "",
        "■ 롤백",
        "kubectl rollout undo",
        "deployment/{svc}",
        "-n rummikub",
    ], l=6.55, t=0.92)

# ─── A-4: 서비스 시작 순서 (의존성) ─────────────────────────────────────────
def slide_startup_order(prs):
    slide = new_slide(prs)
    hdr(slide, "A-4  서비스 시작 순서 — 의존성 다이어그램",
        "DB → Cache → LLM → Backend → Frontend 순서 준수 (역순 종료)")

    # ─ 시작 순서 플로우 ─
    steps = [
        (1, "PostgreSQL\n(DB 레이어)", ":5432",  THEME['blue'],    "helm install postgres\n대기 30초 + wait pod ready"),
        (2, "Redis\n(캐시 레이어)",    ":6379",  THEME['blue'],    "helm install redis\nwait pod ready"),
        (3, "Ollama\n(로컬 LLM)",      ":11434", RGBColor(180,100,0), "helm install ollama\n대기 60초 + 모델 pull"),
        (4, "Game Server\n(백엔드)",   ":8080",  THEME['dark'],    "helm install game-server\nAutoMigrate DB"),
        (5, "AI Adapter\n(LLM 브릿지)",":3001",  THEME['primary'], "helm install ai-adapter\nLLM 키 확인"),
        (6, "Frontend\n(웹 UI)",       ":3000",  THEME['primary'], "helm install frontend\n브라우저 접속 가능"),
        (7, "Admin\n(관리 UI)",        ":3001",  THEME['dark'],    "helm install admin\n관리자 접속 가능"),
    ]

    box_x = [0.28, 1.62, 2.96, 4.3, 5.64, 6.98, 8.32]
    for i, ((num, name, port, color, cmd), x) in enumerate(zip(steps, box_x)):
        y_box = 1.2
        # 순서 원
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.3), Inches(y_box-0.12),
                                   Inches(0.6), Inches(0.5))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.color.rgb = color
        add_tb(slide, x+0.3, y_box-0.1, 0.6, 0.46, str(num), 14, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        # 박스
        add_round(slide, x+0.04, y_box+0.42, 1.16, 1.0, color, line=color)
        add_tb(slide, x+0.04, y_box+0.48, 1.16, 0.52, name, 10.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.04, y_box+1.0, 1.16, 0.35, port, 10, THEME['light'], align=PP_ALIGN.CENTER)
        # 명령어
        add_tb(slide, x+0.04, y_box+1.55, 1.16, 0.7, cmd, 8.5, THEME['text'], align=PP_ALIGN.CENTER)
        # 화살표
        if i < 6:
            add_arrow(slide, x+1.2, y_box+0.92, x+1.34, y_box+0.92, color, Pt(2))

    # 의존성 다이어그램 (하단)
    add_round(slide, 0.28, 3.75, 9.45, 3.35, THEME['light'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.45, 3.8, 9.0, 0.35, "의존성 관계도", 13, THEME['dark'], bold=True)

    # 의존 박스들
    deps = [
        ("PostgreSQL", 1.0, 4.4, THEME['blue']),
        ("Redis",      3.5, 4.4, THEME['blue']),
        ("Ollama",     6.0, 4.4, RGBColor(180,100,0)),
        ("Game Server", 2.25, 5.45, THEME['dark']),
        ("AI Adapter",  5.25, 5.45, THEME['primary']),
        ("Frontend",    1.4, 6.45, THEME['primary']),
        ("Admin",       4.4, 6.45, THEME['dark']),
    ]
    for name, x, y, color in deps:
        add_round(slide, x, y, 1.6, 0.6, color, line=color)
        add_tb(slide, x, y+0.1, 1.6, 0.4, name, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 의존 화살표
    add_arrow(slide, 2.05, 4.7, 2.9, 5.45, THEME['blue'], Pt(1.5))   # PG → GS
    add_arrow(slide, 4.3, 4.7, 3.5, 5.45, THEME['blue'], Pt(1.5))    # Redis → GS
    add_arrow(slide, 4.3, 4.7, 5.85, 5.45, THEME['blue'], Pt(1.5))   # Redis → AI
    add_arrow(slide, 7.6, 4.7, 6.45, 5.45, RGBColor(180,100,0), Pt(1.5))  # Ollama → AI
    add_arrow(slide, 3.05, 6.05, 2.2, 6.45, THEME['dark'], Pt(1.5))   # GS → FE
    add_arrow(slide, 3.05, 6.05, 5.2, 6.45, THEME['dark'], Pt(1.5))   # GS → Admin
    add_arrow(slide, 6.05, 6.05, 5.2, 6.45, THEME['primary'], Pt(1.5)) # AI → Admin

# ─── A-5: 시크릿 관리 구조 ───────────────────────────────────────────────────
def slide_secrets(prs):
    slide = new_slide(prs)
    hdr(slide, "A-5  시크릿 관리 구조",
        "5개 K8s Secret — 이관 시 반드시 백업 후 신규 환경에 재등록")

    secrets = [
        ("google-oauth-secret",  THEME['blue'],
         ["GOOGLE_CLIENT_ID", "GOOGLE_CLIENT_SECRET"],
         "Google Cloud Console\n→ APIs & Services\n→ OAuth 2.0 Client IDs",
         ["console.cloud.google.com", "새 환경 redirect URI 등록 필요",
          "http://localhost:30000/api/auth/callback/google"]),
        ("nextauth-secret",      THEME['primary'],
         ["NEXTAUTH_SECRET", "NEXTAUTH_URL"],
         "openssl rand -base64 32\n→ 랜덤 생성 사용",
         ["각 환경마다 새로 생성", "NEXTAUTH_URL = 실제 접속 URL",
          "예: http://localhost:30000"]),
        ("llm-api-keys",         THEME['accent'],
         ["OPENAI_API_KEY", "ANTHROPIC_API_KEY", "DEEPSEEK_API_KEY"],
         "각 플랫폼에서 발급\n(선택: 최소 1개 필수)",
         ["OpenAI: platform.openai.com", "Anthropic: console.anthropic.com",
          "DeepSeek: platform.deepseek.com"]),
        ("db-secret",            THEME['dark'],
         ["POSTGRES_PASSWORD"],
         "직접 설정\n(복잡한 패스워드 권장)",
         ["이관 후 DB 비밀번호 동일 유지", "또는 dump+restore 후 재설정",
          "PVC 이관 시 동일 값 사용"]),
        ("kakao-secret",         THEME['gray'],
         ["KAKAO_REST_API_KEY"],
         "카카오 개발자 센터\n(선택 — 알림 기능용)",
         ["없어도 시스템 동작", "카카오톡 알림만 비활성화",
          "developers.kakao.com"]),
    ]
    for i, (name, color, keys, source, notes) in enumerate(secrets):
        y = 1.0 + i * 1.22
        # 시크릿명
        add_round(slide, 0.25, y+0.05, 2.35, 1.0, color, line=color)
        add_tb(slide, 0.32, y+0.1, 2.22, 0.4, name, 10.5, THEME['white'], bold=True)
        for j, k in enumerate(keys):
            add_tb(slide, 0.38, y+0.52+j*0.2, 2.1, 0.2, "• "+k, 8.5, THEME['white'])
        # 발급처
        add_round(slide, 2.75, y+0.05, 2.5, 1.0, THEME['light'], line=color, line_w=Pt(1.5))
        add_tb(slide, 2.85, y+0.1, 2.3, 0.22, "발급처", 9, THEME['gray'], bold=True)
        add_tb(slide, 2.85, y+0.32, 2.3, 0.68, source, 10, THEME['text'])
        # 주의사항
        add_round(slide, 5.4, y+0.05, 3.1, 1.0, THEME['light'], line=THEME['gray'], line_w=Pt(1))
        add_tb(slide, 5.5, y+0.1, 2.9, 0.22, "이관 주의사항", 9, THEME['gray'], bold=True)
        for j, note in enumerate(notes):
            add_tb(slide, 5.5, y+0.32+j*0.2, 2.9, 0.2, "• "+note, 8.5, THEME['text'])

    # 백업 명령어
    add_round(slide, 0.25, 7.15, 9.5, 0.45, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.45, 7.18, 9.2, 0.35,
           "백업: kubectl get secret {이름} -n rummikub -o yaml > backup-{이름}.yaml  ⚠️ 평문 저장 — 보안 주의",
           11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── A-6: 일상 운영 절차 ─────────────────────────────────────────────────────
def slide_daily_ops(prs):
    slide = new_slide(prs)
    hdr(slide, "A-6  일상 운영 절차",
        "일일 점검 5항목 · 주간 유지보수 · API 비용 모니터링")

    # 일일 점검
    add_round(slide, 0.25, 0.95, 6.1, 3.6, THEME['light'], line=THEME['primary'], line_w=Pt(2))
    add_rect(slide, 0.25, 0.95, 6.1, 0.5, THEME['primary'])
    add_tb(slide, 0.4, 0.98, 5.8, 0.44, "매일 점검 (Daily Check)", 14, THEME['white'], bold=True)
    daily = [
        ("1", "전체 Pod 상태",      "kubectl get pods -n rummikub",                              THEME['primary']),
        ("2", "에러 로그 확인",     "kubectl logs deployment/game-server --since=1h | grep error",THEME['red']),
        ("3", "AI Adapter 타임아웃","kubectl logs deployment/ai-adapter --since=1h | grep timeout",THEME['red']),
        ("4", "Redis 메모리",       "redis-cli INFO memory | grep used_memory_human",             THEME['blue']),
        ("5", "API 비용 현황",      "curl localhost:30080/admin/stats/ai",                        THEME['accent']),
    ]
    for i, (num, label, cmd, color) in enumerate(daily):
        y = 1.55 + i * 0.58
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.38), Inches(y), Inches(0.38), Inches(0.35))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.color.rgb = color
        add_tb(slide, 0.38, y, 0.38, 0.35, num, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, 0.85, y, 1.4, 0.35, label, 12, THEME['dark'], bold=True)
        add_tb(slide, 0.85, y+0.22, 5.3, 0.22, cmd, 9, THEME['gray'], italic=True)

    # 주간 유지보수
    add_round(slide, 6.5, 0.95, 3.25, 3.6, THEME['light'], line=THEME['dark'], line_w=Pt(2))
    add_rect(slide, 6.5, 0.95, 3.25, 0.5, THEME['dark'])
    add_tb(slide, 6.65, 0.98, 3.0, 0.44, "주간 유지보수 (Weekly)", 14, THEME['white'], bold=True)
    weekly = [
        "PostgreSQL 백업 (pg_dump)",
        "30일 이상 완료된 게임 삭제",
        "Docker 이미지 정리 (prune)",
        "Trivy 보안 스캔 결과 확인",
        "LLM API 잔액 확인",
        "K8s 이벤트 로그 검토",
    ]
    for i, item in enumerate(weekly):
        y = 1.55 + i * 0.5
        add_tb(slide, 6.65, y, 3.0, 0.42, "✓  " + item, 11.5, THEME['text'])

    # 비용 모니터링 다이어그램
    add_round(slide, 0.25, 4.65, 9.5, 2.5, THEME['light'], line=THEME['accent'], line_w=Pt(2))
    add_rect(slide, 0.25, 4.65, 9.5, 0.5, THEME['accent'])
    add_tb(slide, 0.4, 4.68, 9.0, 0.44, "API 비용 모니터링 구조", 14, THEME['white'], bold=True)

    # 비용 플로우
    nodes = [
        ("AI Adapter\n비용 집계",  1.0,  5.35, THEME['primary']),
        ("일일 한도\n$20",        3.2,  5.35, THEME['dark']),
        ("사용자 한도\n$5/시간",  5.4,  5.35, THEME['dark']),
        ("초과 시\n요청 차단",    7.6,  5.35, THEME['red']),
    ]
    for label, x, y, color in nodes:
        add_round(slide, x, y, 1.8, 0.85, color, line=color)
        add_tb(slide, x, y+0.12, 1.8, 0.62, label, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    for x in [2.8, 5.0, 7.2]:
        add_arrow(slide, x, 5.77, x+0.4, 5.77, THEME['accent'], Pt(2))

    add_tb(slide, 0.45, 6.3, 9.1, 0.75,
           "모델별 비용: DeepSeek $0.039/게임  |  OpenAI GPT-5-mini 가변  |  Claude Sonnet 4 $1.11/게임  |  Ollama 무료\n"
           "Admin API: GET /admin/stats/ai → today_cost_usd / daily_limit_usd / model_breakdown",
           11, THEME['text'])

# ─── A-7: 장애 대응 플로우 ───────────────────────────────────────────────────
def slide_incident_flow(prs):
    slide = new_slide(prs)
    hdr(slide, "A-7  장애 대응 플로우차트",
        "장애 감지 → 원인 분류 → 유형별 대응 → 복구 확인 → 사후 기록")

    # 플로우 노드
    def flow_rect(slide, x, y, w, h, text, color, font_size=12):
        add_round(slide, x, y, w, h, color, line=color)
        add_tb(slide, x, y+h*0.2, w, h*0.65, text, font_size, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    def flow_diamond(slide, x, y, w, h, text, color, font_size=11):
        add_diam(slide, x, y, w, h, color, line=color)
        add_tb(slide, x+w*0.1, y+h*0.25, w*0.8, h*0.55, text, font_size, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 시작
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.1), Inches(0.98), Inches(1.8), Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb = THEME['dark']; s.line.color.rgb = THEME['dark']
    add_tb(slide, 4.1, 1.02, 1.8, 0.47, "장애 감지", 13, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 5.0, 1.53, 5.0, 1.85, THEME['dark'], Pt(2))

    # 감지 방법
    flow_rect(slide, 3.4, 1.85, 3.2, 0.62, "Pod 상태 확인\nkubectl get pods -n rummikub", THEME['primary'], 10)
    add_arrow(slide, 5.0, 2.47, 5.0, 2.75, THEME['dark'], Pt(2))

    # 분기: Pod 상태?
    flow_diamond(slide, 3.7, 2.75, 2.6, 0.85, "Pod\nRunning?", THEME['accent'])
    # Yes → 로그 확인
    add_arrow(slide, 6.3, 3.18, 7.5, 3.18, THEME['primary'], Pt(2))
    add_tb(slide, 6.35, 3.0, 1.0, 0.25, "Running", 9, THEME['primary'], bold=True)
    flow_rect(slide, 7.5, 2.85, 2.0, 0.65, "로그 분석\n(에러 패턴 확인)", THEME['primary'], 10)

    # No → 재시작
    add_arrow(slide, 5.0, 3.6, 5.0, 3.9, THEME['red'], Pt(2))
    add_tb(slide, 5.1, 3.65, 0.8, 0.25, "Not Ready", 9, THEME['red'], bold=True)
    flow_rect(slide, 3.55, 3.9, 2.9, 0.65, "kubectl rollout restart\ndeployment/{svc}", THEME['red'], 10)

    # 분기: 원인 분류
    add_arrow(slide, 7.5, 3.18, 7.5, 4.05, THEME['primary'], Pt(2))
    add_arrow(slide, 5.0, 4.55, 5.0, 4.85, THEME['red'], Pt(2))

    flow_diamond(slide, 3.6, 4.85, 2.8, 0.9, "서비스\n종류?", THEME['dark'])

    # 분기 4방향
    branches = [
        (0.6, 4.95, "frontend", THEME['primary']),
        (0.6, 5.85, "game-server", THEME['dark']),
        (6.5, 4.95, "ai-adapter", THEME['primary']),
        (6.5, 5.85, "DB/Redis", THEME['blue']),
    ]
    for x, y, label, color in branches:
        add_round(slide, x, y, 1.9, 0.58, color, line=color)
        add_tb(slide, x, y+0.1, 1.9, 0.38, label, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    add_arrow(slide, 3.6, 5.3, 2.5, 5.24, THEME['primary'], Pt(1.5))
    add_arrow(slide, 3.6, 5.3, 2.5, 6.14, THEME['dark'], Pt(1.5))
    add_arrow(slide, 6.4, 5.3, 6.5, 5.24, THEME['primary'], Pt(1.5))
    add_arrow(slide, 6.4, 5.3, 6.5, 6.14, THEME['blue'], Pt(1.5))

    add_tb(slide, 0.3, 5.55, 2.5, 0.25, "→ A-8 슬라이드 참조", 9, THEME['primary'], italic=True)
    add_tb(slide, 0.3, 6.5, 2.5, 0.25, "→ A-8 슬라이드 참조", 9, THEME['dark'], italic=True)
    add_tb(slide, 6.5, 5.55, 2.5, 0.25, "→ A-8 슬라이드 참조", 9, THEME['primary'], italic=True)
    add_tb(slide, 6.5, 6.5, 2.5, 0.25, "→ A-9 슬라이드 참조", 9, THEME['blue'], italic=True)

    # 복구 확인
    flow_rect(slide, 3.4, 6.55, 3.2, 0.58, "복구 확인: curl :30080/health\n+ kubectl get pods", THEME['primary'], 10)
    add_arrow(slide, 5.0, 6.55, 5.0, 6.43, THEME['dark'], Pt(1.5))  # 위에서 내려옴 없이 독립

    # 사후 기록
    add_round(slide, 3.4, 7.2, 3.2, 0.38, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 3.4, 7.23, 3.2, 0.3, "사후 기록: work_logs/incidents/", 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── A-8: 장애 유형별 대응 (1/2) ─────────────────────────────────────────────
def slide_incident_1(prs):
    slide = new_slide(prs)
    hdr(slide, "A-8  장애 유형별 대응 (1/2)",
        "frontend 접속 불가 · game-server 오류 · ai-adapter 응답 없음")

    cases = [
        {
            "title": "1. frontend 접속 불가 (http://localhost:30000)",
            "color": THEME['primary'],
            "symptoms": ["브라우저 ERR_CONNECTION_REFUSED", "흰 화면 또는 로딩 무한 대기"],
            "steps": [
                "kubectl get pod -n rummikub -l app=frontend",
                "kubectl logs -n rummikub deployment/frontend --tail=50",
                "kubectl describe pod -n rummikub -l app=frontend",
                "kubectl rollout restart deployment/frontend -n rummikub",
                "kubectl rollout status deployment/frontend -n rummikub  # 완료 대기",
            ],
            "note": "⚠ NEXTAUTH_URL이 실제 접속 URL과 다르면 OAuth 로그인 실패 발생"
        },
        {
            "title": "2. game-server API 실패 / WebSocket 연결 불가",
            "color": THEME['dark'],
            "symptoms": ["API 호출 시 502/503 응답", "WebSocket onclose 즉시 발생"],
            "steps": [
                "kubectl describe deployment game-server -n rummikub",
                "kubectl exec -n rummikub deployment/game-server -- env | grep -E 'DB_|REDIS_'",
                "kubectl exec -n rummikub deployment/game-server -- curl localhost:8080/health",
                "kubectl rollout restart deployment/game-server -n rummikub",
                "kubectl logs -n rummikub deployment/game-server --since=5m | grep -i error",
            ],
            "note": "⚠ DB 연결 실패 시 game-server 재시작 전 postgres/redis 상태 먼저 확인"
        },
        {
            "title": "3. ai-adapter 응답 없음 (AI 차례 멈춤)",
            "color": THEME['primary'],
            "symptoms": ["AI 차례에서 1000초 대기 후 강제 드로우", "timeout 로그 반복"],
            "steps": [
                "kubectl logs -n rummikub deployment/ai-adapter --tail=100 | grep -E 'error|timeout|INVALID'",
                "kubectl get secret llm-api-keys -n rummikub -o jsonpath='{.data.OPENAI_API_KEY}' | base64 -d",
                "kubectl exec -n rummikub deployment/ollama -- ollama list  # Ollama 모델 확인",
                "kubectl rollout restart deployment/ai-adapter -n rummikub",
                "# LLM API 키 만료 시: kubectl apply 로 시크릿 재등록 후 재시작",
            ],
            "note": "⚠ Ollama가 응답 없으면 ollama 재시작 → qwen2.5:3b pull 재실행 필요"
        },
    ]

    for i, case in enumerate(cases):
        y = 1.0 + i * 2.08
        color = case['color']
        add_round(slide, 0.25, y, 9.5, 1.95, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, 0.25, y, 9.5, 0.45, color)
        add_tb(slide, 0.38, y+0.06, 9.1, 0.36, case['title'], 12.5, THEME['white'], bold=True)
        # 증상
        add_tb(slide, 0.38, y+0.52, 1.1, 0.22, "증상:", 9.5, color, bold=True)
        for j, s in enumerate(case['symptoms']):
            add_tb(slide, 0.38, y+0.72+j*0.2, 3.5, 0.22, "• "+s, 10, THEME['text'])
        # 명령어
        add_tb(slide, 3.85, y+0.52, 1.1, 0.22, "대응 순서:", 9.5, color, bold=True)
        for j, cmd in enumerate(case['steps']):
            add_tb(slide, 3.85, y+0.72+j*0.22, 5.7, 0.22, f"  {j+1}. "+cmd, 8.5, THEME['text'])
        # 주의
        add_tb(slide, 0.38, y+1.75, 9.2, 0.2, case['note'], 9, THEME['red'], italic=True)

# ─── A-9: 장애 유형별 대응 (2/2) ─────────────────────────────────────────────
def slide_incident_2(prs):
    slide = new_slide(prs)
    hdr(slide, "A-9  장애 유형별 대응 (2/2)",
        "PostgreSQL 연결 실패 · Redis 장애 · 전체 시스템 재시작 순서")

    cases = [
        {
            "title": "4. PostgreSQL 연결 실패 (로그인 불가 · 방 목록 오류)",
            "color": THEME['blue'],
            "steps": [
                "kubectl get pod -n rummikub -l app=postgres",
                "kubectl exec -n rummikub deployment/postgres -- psql -U rummikub -c 'SELECT 1;'",
                "kubectl get pvc -n rummikub  # PVC 마운트 상태 확인",
                "kubectl rollout restart deployment/postgres -n rummikub  # PVC 보호로 데이터 유지",
            ],
            "note": "⚠ PVC 삭제하면 DB 데이터 전체 손실. 재시작만으로 복구 안 되면 백업 복원 절차 진행"
        },
        {
            "title": "5. Redis 장애 (게임 상태 손실 · 타이머 오작동)",
            "color": THEME['blue'],
            "steps": [
                "kubectl exec -n rummikub deployment/redis -- redis-cli PING  # PONG 확인",
                "kubectl exec -n rummikub deployment/redis -- redis-cli DBSIZE  # 키 수 확인",
                "kubectl rollout restart deployment/redis -n rummikub",
                "# 주의: 재시작 시 진행 중인 게임 상태 손실. 사용자에게 사전 공지 필요",
            ],
            "note": "⚠ Redis 재시작 시 현재 진행 중인 게임은 모두 중단됨. 영속 데이터(DB)는 영향 없음"
        },
    ]
    for i, case in enumerate(cases):
        y = 1.0 + i * 2.0
        color = case['color']
        add_round(slide, 0.25, y, 9.5, 1.88, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, 0.25, y, 9.5, 0.45, color)
        add_tb(slide, 0.38, y+0.06, 9.1, 0.36, case['title'], 12.5, THEME['white'], bold=True)
        add_tb(slide, 0.38, y+0.52, 1.1, 0.22, "대응 순서:", 9.5, color, bold=True)
        for j, cmd in enumerate(case['steps']):
            add_tb(slide, 0.38, y+0.72+j*0.22, 9.1, 0.22, f"  {j+1}. "+cmd, 9.5, THEME['text'])
        add_tb(slide, 0.38, y+1.68, 9.2, 0.2, case['note'], 9, THEME['red'], italic=True)

    # 전체 재시작 순서 (도형 플로우)
    add_round(slide, 0.25, 5.1, 9.5, 2.65, THEME['light'], line=THEME['dark'], line_w=Pt(2))
    add_rect(slide, 0.25, 5.1, 9.5, 0.45, THEME['dark'])
    add_tb(slide, 0.38, 5.13, 9.1, 0.38, "6. 전체 시스템 재시작 순서 (의존성 순서 준수)", 13, THEME['white'], bold=True)

    restart_order = [
        ("postgres\n(30초 대기)", THEME['blue']),
        ("redis\n(5초 대기)",     THEME['blue']),
        ("ollama\n(5초 대기)",    RGBColor(180,100,0)),
        ("game-server\n(5초)",    THEME['dark']),
        ("ai-adapter\n(5초)",     THEME['primary']),
        ("frontend\nadmin",       THEME['primary']),
    ]
    for i, (label, color) in enumerate(restart_order):
        x = 0.45 + i * 1.57
        add_round(slide, x, 5.65, 1.42, 0.88, color, line=color)
        add_tb(slide, x, 5.72, 1.42, 0.74, label, 10.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        if i < 5:
            add_arrow(slide, x+1.42, 6.09, x+1.57, 6.09, THEME['accent'], Pt(2.5))

    add_tb(slide, 0.38, 6.62, 9.2, 0.35,
           "kubectl rollout restart deployment/{이름} -n rummikub  →  각 단계 완료 후 다음 진행",
           11, THEME['text'], italic=True)
    add_tb(slide, 0.38, 6.97, 9.2, 0.35,
           "최종 확인: kubectl get pods -n rummikub  →  모든 Pod STATUS=Running, READY=1/1",
           11, THEME['dark'], bold=True)

# ─── A-10: 배포 절차 ─────────────────────────────────────────────────────────
def slide_deploy(prs):
    slide = new_slide(prs)
    hdr(slide, "A-10  배포 절차 — GitOps · 수동 · 롤백",
        "정상: GitOps(권장) | 긴급: 수동 Helm | 문제: kubectl rollout undo")

    # 3개 트랙
    tracks = [
        ("GitOps 배포 (정상 경로)", THEME['primary'], [
            "1. git push origin main",
            "2. GitLab CI 17단계 자동 실행",
            "   (lint → build → scan → deploy → verify)",
            "3. CI 통과 시 ArgoCD 자동 Sync",
            "4. K8s Rolling Update 실행",
            "5. 헬스체크 통과 → 배포 완료",
            "",
            "확인: kubectl get application -n argocd",
            "또는 ArgoCD UI: http://localhost:30087",
        ], "권장 경로. 테스트 검증 후 배포"),
        ("수동 Helm 배포 (긴급 시)", THEME['accent'], [
            "# values.yaml 이미지 태그 수정 후:",
            "helm upgrade {svc} helm/charts/{svc} \\",
            "  -n rummikub \\",
            "  --set image.tag=새태그",
            "",
            "# 배포 진행 확인",
            "kubectl rollout status \\",
            "  deployment/{svc} -n rummikub",
            "",
            "※ CI/CD 우회 — 보안 스캔 미실행",
        ], "CI 없이 직접 배포. 긴급 패치 시만 사용"),
        ("롤백 (문제 발생 시)", THEME['red'], [
            "# 직전 버전으로 즉시 롤백",
            "kubectl rollout undo \\",
            "  deployment/{svc} -n rummikub",
            "",
            "# 히스토리 확인 후 특정 버전 롤백",
            "kubectl rollout history \\",
            "  deployment/{svc} -n rummikub",
            "kubectl rollout undo \\",
            "  deployment/{svc} --to-revision=3",
        ], "배포 후 문제 즉시 발견 시. ArgoCD sync 일시 정지 필요"),
    ]
    for i, (title, color, steps, note) in enumerate(tracks):
        x = 0.25 + i * 3.25
        add_round(slide, x, 0.95, 3.1, 5.95, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, x, 0.95, 3.1, 0.5, color)
        add_tb(slide, x+0.1, 0.98, 2.9, 0.44, title, 12, THEME['white'], bold=True)
        tb = slide.shapes.add_textbox(Inches(x+0.15), Inches(1.52), Inches(2.8), Inches(4.8))
        tf = tb.text_frame; tf.word_wrap = True
        for j, line in enumerate(steps):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(10.5 if not line.startswith("#") else 9.5)
            p.font.color.rgb = THEME['text'] if not line.startswith("#") else THEME['gray']
            p.font.italic = line.startswith("#") or line.startswith("※")
            p.font.name = FONT; p.space_after = Pt(3)
        add_round(slide, x+0.1, 6.08, 2.9, 0.72, color, line=color)
        add_tb(slide, x+0.15, 6.12, 2.8, 0.62, note, 10, THEME['white'], align=PP_ALIGN.CENTER)

# ─── A-11: 스케일링 전략 ─────────────────────────────────────────────────────
def slide_scaling(prs):
    slide = new_slide(prs)
    hdr(slide, "A-11  스케일링 전략 — 현재 vs 확장",
        "동시접속 100~200명 달성을 위한 단계적 확장 방안")

    # 현재 상태
    add_round(slide, 0.25, 0.95, 4.55, 3.3, THEME['light'], line=THEME['gray'], line_w=Pt(2))
    add_rect(slide, 0.25, 0.95, 4.55, 0.48, THEME['gray'])
    add_tb(slide, 0.38, 0.98, 4.2, 0.4, "현재 (단일 노드 · 즉시 가능)", 13, THEME['white'], bold=True)
    current = [
        ("game-server",  "×1 replica",   THEME['dark']),
        ("ai-adapter",   "×1 replica",   THEME['primary']),
        ("frontend",     "×1 replica",   THEME['primary']),
        ("postgres",     "×1 (단일)",    THEME['blue']),
        ("redis",        "×1 (단일)",    THEME['blue']),
        ("ollama",       "CPU (25s/턴)", RGBColor(180,100,0)),
    ]
    for i, (svc, rep, color) in enumerate(current):
        y = 1.52 + i * 0.46
        add_round(slide, 0.38, y, 1.8, 0.38, color, line=color)
        add_tb(slide, 0.4, y+0.06, 1.76, 0.28, svc, 11, THEME['white'], bold=True)
        add_tb(slide, 2.28, y+0.06, 2.3, 0.28, rep, 11, THEME['text'])

    # 화살표
    add_arrow(slide, 4.8, 2.6, 5.45, 2.6, THEME['accent'], Pt(3))
    add_tb(slide, 4.72, 2.3, 0.9, 0.3, "확장", 11, THEME['accent'], bold=True, align=PP_ALIGN.CENTER)

    # 목표 상태
    add_round(slide, 5.45, 0.95, 4.3, 3.3, THEME['light'], line=THEME['primary'], line_w=Pt(2))
    add_rect(slide, 5.45, 0.95, 4.3, 0.48, THEME['primary'])
    add_tb(slide, 5.58, 0.98, 4.0, 0.4, "목표 (100~200명 지원)", 13, THEME['white'], bold=True)
    target = [
        ("game-server",  "×3 replica  (stateless 수평확장)",   THEME['dark']),
        ("ai-adapter",   "×2 replica  (동시 커넥션 2배)",       THEME['primary']),
        ("frontend",     "×2 replica  (부하 분산)",             THEME['primary']),
        ("postgres",     "×1 (충분) / Cloud SQL 옵션",          THEME['blue']),
        ("redis",        "×1 (충분) / Redis Cluster 옵션",      THEME['blue']),
        ("ollama",       "GPU 노드  (3s/턴 · 10배 처리)",       THEME['accent']),
    ]
    for i, (svc, rep, color) in enumerate(target):
        y = 1.52 + i * 0.46
        add_round(slide, 5.58, y, 1.8, 0.38, color, line=color)
        add_tb(slide, 5.6, y+0.06, 1.76, 0.28, svc, 11, THEME['white'], bold=True)
        add_tb(slide, 7.48, y+0.06, 2.15, 0.28, rep, 10, THEME['text'])

    # 단계별 확장 명령어
    add_round(slide, 0.25, 4.35, 9.5, 2.85, THEME['light'], line=THEME['dark'], line_w=Pt(1.5))
    add_rect(slide, 0.25, 4.35, 9.5, 0.45, THEME['dark'])
    add_tb(slide, 0.38, 4.38, 9.1, 0.38, "즉시 적용 가능한 명령어 (추가 노드 없이)", 13, THEME['white'], bold=True)

    cmds = [
        ("Step 1", "game-server 3배 확장 (즉시 가능 — stateless)",
         "kubectl scale deployment game-server -n rummikub --replicas=3"),
        ("Step 2", "ai-adapter 2배 확장",
         "kubectl scale deployment ai-adapter -n rummikub --replicas=2"),
        ("Step 3", "Ollama CPU 리소스 상향",
         "helm upgrade ollama helm/charts/ollama -n rummikub --set resources.limits.cpu=8"),
    ]
    for i, (step, desc, cmd) in enumerate(cmds):
        y = 4.9 + i * 0.74
        s = add_round(slide, 0.38, y, 0.85, 0.44, THEME['primary'], line=THEME['primary'])
        add_tb(slide, 0.38, y+0.1, 0.85, 0.28, step, 10, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        add_tb(slide, 1.35, y+0.02, 3.0, 0.22, desc, 11, THEME['dark'], bold=True)
        add_tb(slide, 1.35, y+0.24, 8.2, 0.22, cmd, 9.5, THEME['blue'], italic=True)

# ─── A-12: 클라우드 이관 아키텍처 (GKE) ─────────────────────────────────────
def slide_cloud(prs):
    slide = new_slide(prs)
    hdr(slide, "A-12  클라우드 이관 아키텍처 — GKE 기준",
        "단일 노드 Docker Desktop K8s → Google Kubernetes Engine (월 $180~280)")

    # 왼쪽: 현재
    add_round(slide, 0.2, 0.95, 4.1, 6.3, THEME['light'], line=THEME['gray'], line_w=Pt(2))
    add_rect(slide, 0.2, 0.95, 4.1, 0.48, THEME['gray'])
    add_tb(slide, 0.32, 0.98, 3.8, 0.38, "현재: Docker Desktop K8s", 12, THEME['white'], bold=True)

    current_nodes = [
        ("Windows 11 Host", THEME['lgray'], 0.32, 1.52),
    ]
    add_round(slide, 0.32, 1.52, 3.82, 5.55, THEME['white'], line=THEME['primary'], line_w=Pt(1.5))
    add_tb(slide, 0.42, 1.55, 3.5, 0.28, "단일 노드 (LG Gram)", 10, THEME['primary'], italic=True)
    services_cur = ["frontend ×1", "admin ×1", "game-server ×1", "ai-adapter ×1", "postgres ×1", "redis ×1", "ollama (CPU) ×1"]
    for i, svc in enumerate(services_cur):
        c = THEME['primary'] if i < 4 else THEME['blue'] if i < 6 else RGBColor(180,100,0)
        add_round(slide, 0.45, 1.88+i*0.72, 3.55, 0.6, c, line=c)
        add_tb(slide, 0.5, 1.94+i*0.72, 3.45, 0.48, svc, 11.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 이관 화살표
    add_arrow(slide, 4.3, 4.1, 5.35, 4.1, THEME['accent'], Pt(4))
    add_tb(slide, 4.18, 3.65, 1.3, 0.42, "이관\n(GKE)", 11, THEME['accent'], bold=True, align=PP_ALIGN.CENTER)

    # 오른쪽: GKE
    add_round(slide, 5.35, 0.95, 4.4, 6.3, THEME['light'], line=THEME['primary'], line_w=Pt(2))
    add_rect(slide, 5.35, 0.95, 4.4, 0.48, THEME['primary'])
    add_tb(slide, 5.48, 0.98, 4.1, 0.38, "목표: GKE Autopilot", 12, THEME['white'], bold=True)

    # General 노드 풀
    add_round(slide, 5.48, 1.52, 4.12, 3.1, THEME['lblue'], line=THEME['blue'], line_w=Pt(1.5))
    add_tb(slide, 5.58, 1.55, 3.8, 0.28, "General Node Pool (n2-std-4 × 2~5)", 9.5, THEME['blue'], italic=True)
    general_svcs = [
        ("frontend ×2", THEME['primary']),
        ("admin ×1", THEME['dark']),
        ("game-server ×3", THEME['dark']),
        ("ai-adapter ×2", THEME['primary']),
    ]
    for i, (svc, color) in enumerate(general_svcs):
        col = i % 2; row = i // 2
        x = 5.62 + col * 2.0; y = 1.9 + row * 0.82
        add_round(slide, x, y, 1.85, 0.65, color, line=color)
        add_tb(slide, x, y+0.1, 1.85, 0.45, svc, 11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # GPU 노드 풀
    add_round(slide, 5.48, 4.72, 4.12, 1.1, THEME['lyellow'], line=THEME['accent'], line_w=Pt(1.5))
    add_tb(slide, 5.58, 4.75, 3.8, 0.28, "GPU Node Pool (n1-std-4 + NVIDIA T4)", 9.5, RGBColor(180,100,0), italic=True)
    add_round(slide, 5.65, 5.05, 3.8, 0.65, RGBColor(180,100,0), line=RGBColor(180,100,0))
    add_tb(slide, 5.65, 5.15, 3.8, 0.45, "ollama ×1  (GPU 추론 · 3s/턴)", 11.5, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 공유 인프라
    add_round(slide, 5.48, 5.92, 4.12, 1.2, THEME['light'], line=THEME['dark'], line_w=Pt(1.5))
    add_tb(slide, 5.58, 5.95, 3.8, 0.28, "공유 인프라", 9.5, THEME['dark'], italic=True)
    shared = [("postgres ×1\n(Cloud SQL 옵션)", THEME['blue']), ("redis ×1\n(Memorystore 옵션)", THEME['blue'])]
    for i, (svc, color) in enumerate(shared):
        x = 5.62 + i * 2.0
        add_round(slide, x, 6.22, 1.85, 0.72, color, line=color)
        add_tb(slide, x, 6.3, 1.85, 0.58, svc, 10, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

    # 주요 변경사항
    add_round(slide, 0.2, 7.17, 9.55, 0.42, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 0.35, 7.2, 9.3, 0.34,
           "주요 변경: NodePort → LoadBalancer  |  StorageClass: standard → standard-rwo  |  도메인+HTTPS 설정 필요  |  월 $180~280",
           11, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ─── A-13: 이관 전 체크리스트 ────────────────────────────────────────────────
def slide_checklist_pre(prs):
    slide = new_slide(prs)
    hdr(slide, "A-13  이관 전 체크리스트",
        "인계 전 반드시 완료해야 하는 항목 — 체크 후 서명")

    categories = [
        ("데이터 백업", THEME['blue'], [
            ("PostgreSQL pg_dump 백업 실행 완료",           True),
            ("백업 파일 크기 확인 (0B 아님)",               True),
            ("Redis 진행 중 게임 없음 확인 (KEYS game:*)",  True),
            ("백업 파일 외부 저장소 업로드",                True),
        ]),
        ("시크릿 이관", THEME['red'], [
            ("google-oauth-secret YAML 추출",               True),
            ("nextauth-secret YAML 추출",                   True),
            ("llm-api-keys YAML 추출",                      True),
            ("db-secret YAML 추출",                         True),
            ("백업 파일 암호화 보관 (평문 주의)",            True),
        ]),
        ("이미지 확보", THEME['primary'], [
            ("DockerHub rummiarena/* 이미지 접근 확인",     True),
            ("로컬 이미지 tar 추출 또는 DockerHub push",    True),
            ("GitLab CI 재빌드 가능 여부 확인",             True),
        ]),
        ("문서 인수인계", THEME['dark'], [
            ("운영자 매뉴얼 전달 (06-operations/10번)",     True),
            ("사용자 매뉴얼 전달 (06-operations/11번)",     True),
            ("운영환경 구성 가이드 전달 (12번)",             True),
            ("시크릿 값 별도 보안 채널 전달",               True),
        ]),
        ("신규 환경 검증", THEME['accent'], [
            ("kubectl get nodes → Ready 확인",              False),
            ("7개 서비스 모두 Running 확인",                False),
            ("curl :30080/health → ok 응답",                False),
            ("브라우저 게임 접속 + AI 대전 1회 테스트",     False),
        ]),
    ]

    col_w = 4.55
    for ci, (cat, color, items) in enumerate(categories):
        col = ci % 2; row = ci // 2
        x = 0.25 + col * 4.85
        y = 0.95 + row * 2.32
        if ci == 4:  # 마지막은 전체 너비
            x = 0.25; col_w_use = 9.5
        else:
            col_w_use = 4.55
        add_round(slide, x, y, col_w_use, 2.18, THEME['light'], line=color, line_w=Pt(2))
        add_rect(slide, x, y, col_w_use, 0.44, color)
        add_tb(slide, x+0.12, y+0.05, col_w_use-0.24, 0.36, cat, 13, THEME['white'], bold=True)
        for j, (item, done) in enumerate(items):
            yy = y + 0.52 + j * 0.38
            check_color = THEME['primary'] if done else THEME['gray']
            check = "☑" if done else "☐"
            add_tb(slide, x+0.15, yy, 0.32, 0.32, check, 14, check_color, bold=True)
            status = "완료" if done else "신규"
            s = add_round(slide, x+0.5, yy+0.04, 0.52, 0.26, check_color, line=check_color)
            add_tb(slide, x+0.5, yy+0.06, 0.52, 0.22, status, 8, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
            add_tb(slide, x+1.1, yy+0.04, col_w_use-1.22, 0.28, item, 11, THEME['text'])

# ─── A-14: 이관 후 검증 절차 ─────────────────────────────────────────────────
def slide_checklist_post(prs):
    slide = new_slide(prs)
    hdr(slide, "A-14  이관 후 검증 절차",
        "신규 환경에서 서비스 정상 동작 최종 확인 — 7단계 검증")

    steps = [
        ("1", "전체 Pod 상태 확인",
         "kubectl get pods -n rummikub",
         "모든 Pod STATUS=Running, READY=1/1",
         THEME['primary']),
        ("2", "게임 서버 헬스체크",
         'curl http://localhost:30080/health',
         '{"status":"ok","redis":"connected","db":"connected"}',
         THEME['dark']),
        ("3", "AI Adapter 헬스체크",
         "curl http://localhost:30081/health",
         '{"status":"UP"}',
         THEME['primary']),
        ("4", "브라우저 게임 접속",
         "http://localhost:30000 접속 → Google OAuth 로그인",
         "로그인 성공 · 로비 화면 표시 · ELO 순위 로드",
         THEME['accent']),
        ("5", "방 생성 + AI 대전 테스트",
         "방 생성 → DeepSeek AI 추가 → 게임 시작 → 5턴 진행",
         "AI가 정상 배치 · 타이머 동작 · 점수 갱신",
         THEME['dark']),
        ("6", "보안 스캔",
         "GitLab CI 파이프라인 실행 → Trivy 스캔",
         "Critical/High CVE = 0건",
         THEME['red']),
        ("7", "DB 백업 테스트",
         "kubectl exec deployment/postgres -- pg_dump -U rummikub rummikub | wc -l",
         "0보다 큰 숫자 → 백업 가능",
         THEME['blue']),
    ]

    for i, (num, title, cmd, expected, color) in enumerate(steps):
        y = 1.0 + i * 0.88
        # 번호
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.25), Inches(y+0.1),
                                   Inches(0.52), Inches(0.52))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.color.rgb = color
        add_tb(slide, 0.25, y+0.1, 0.52, 0.52, num, 15, THEME['white'], bold=True, align=PP_ALIGN.CENTER)
        # 제목
        add_tb(slide, 0.88, y+0.1, 2.5, 0.28, title, 13, color, bold=True)
        # 명령어
        add_round(slide, 0.88, y+0.4, 4.4, 0.3, THEME['light'], line=THEME['gray'], line_w=Pt(1))
        add_tb(slide, 0.97, y+0.43, 4.25, 0.24, cmd, 9, THEME['blue'], italic=True)
        # 기대 결과
        add_round(slide, 5.45, y+0.4, 4.3, 0.3, THEME['light'], line=color, line_w=Pt(1))
        add_tb(slide, 5.55, y+0.43, 4.12, 0.24, "→ " + expected, 9, THEME['dark'])

    # 열 헤더
    add_round(slide, 0.88, 0.93, 4.4, 0.42, THEME['primary'], line=THEME['primary'])
    add_tb(slide, 0.97, 0.95, 4.25, 0.36, "검증 명령어 / 절차", 12, THEME['white'], bold=True)
    add_round(slide, 5.45, 0.93, 4.3, 0.42, THEME['dark'], line=THEME['dark'])
    add_tb(slide, 5.55, 0.95, 4.12, 0.36, "기대 결과 (PASS 기준)", 12, THEME['white'], bold=True)

    # 합격 기준
    add_round(slide, 0.25, 7.17, 9.5, 0.42, THEME['primary'], line=THEME['primary'])
    add_tb(slide, 0.38, 7.2, 9.2, 0.34,
           "7단계 모두 PASS → 이관 완료 선언  |  FAIL 항목 발생 시 A-8/A-9 장애 대응 절차 참조",
           12, THEME['white'], bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 메인: 기존 PPTX에 append
# ══════════════════════════════════════════════════════════════════════════════
def main():
    base_path = os.path.join(os.path.dirname(__file__), "RummiArena_프로젝트종료보고서.pptx")
    prs = Presentation(base_path)

    print(f"기존 슬라이드 수: {len(prs.slides)}")
    print("별첨 슬라이드 추가 중...")

    slide_annex_cover(prs);    print("  A-0  별첨 표지")
    slide_env_arch(prs);       print("  A-1  현재 운영 환경 아키텍처")
    slide_k8s_detail(prs);     print("  A-2  K8s 서비스 구조도")
    slide_gitops(prs);         print("  A-3  GitOps 배포 파이프라인")
    slide_startup_order(prs);  print("  A-4  서비스 시작 순서")
    slide_secrets(prs);        print("  A-5  시크릿 관리 구조")
    slide_daily_ops(prs);      print("  A-6  일상 운영 절차")
    slide_incident_flow(prs);  print("  A-7  장애 대응 플로우")
    slide_incident_1(prs);     print("  A-8  장애 유형별 대응 (1/2)")
    slide_incident_2(prs);     print("  A-9  장애 유형별 대응 (2/2)")
    slide_deploy(prs);         print("  A-10 배포 절차")
    slide_scaling(prs);        print("  A-11 스케일링 전략")
    slide_cloud(prs);          print("  A-12 클라우드 이관 아키텍처")
    slide_checklist_pre(prs);  print("  A-13 이관 전 체크리스트")
    slide_checklist_post(prs); print("  A-14 이관 후 검증 절차")

    prs.save(base_path)
    print(f"\n완료: {base_path}")
    print(f"총 슬라이드: {len(prs.slides)}장 (본문 20 + 별첨 15)")

if __name__ == "__main__":
    main()
